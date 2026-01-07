#!/usr/bin/env python3
"""
Roadmap Synthesis Tool - Streamlit Web Interface
"""

import streamlit as st
from pathlib import Path
import os
from datetime import datetime
from typing import Optional, List, Dict
import pandas as pd
from enum import Enum
from dataclasses import dataclass

# Import functions from roadmap.py
from roadmap import (
    parse_document, chunk_text, chunk_with_fallback, index_chunks, retrieve_chunks,
    generate_roadmap, format_for_persona, init_db,
    VALID_LENSES, OUTPUT_DIR, DATA_DIR, MATERIALS_DIR,
    validate_api_keys, ContextGraph, generate_embeddings,
    load_questions, save_questions, load_answers, save_answers,
    load_decisions, save_decisions,
    load_architecture_documents, scan_architecture_documents, generate_architecture_alignment,
    parse_roadmap_for_analysis, extract_engineering_questions_from_alignment,
    add_architecture_questions_to_system, save_alignment_analysis,
    load_alignment_analysis, format_alignment_report,
    load_competitor_developments, add_competitor_development, get_competitor_development,
    load_analyst_assessments, generate_analyst_assessment, format_analyst_assessment_markdown,
    UnifiedContextGraph, sync_all_to_graph, retrieve_with_authority, AUTHORITY_LEVELS
)

# Page configuration
st.set_page_config(
    page_title="Roadmap Synth",
    page_icon="🗺️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'index_stats' not in st.session_state:
    st.session_state.index_stats = None


# ========== UTILITY FUNCTIONS ==========

def get_index_stats() -> Optional[Dict]:
    """Get statistics about indexed materials"""
    try:
        db = init_db()
        table = db.open_table("roadmap_chunks")

        # Get all chunks
        chunks = table.to_pandas()

        if chunks.empty:
            return None

        # Calculate statistics
        total_chunks = len(chunks)
        total_tokens = chunks['token_count'].sum()

        # Breakdown by lens
        lens_counts = chunks['lens'].value_counts().to_dict()

        # Recent sources
        recent_sources = chunks.nlargest(10, 'created_at')[['source_file', 'created_at', 'lens']].to_dict('records')

        return {
            'total_chunks': total_chunks,
            'total_tokens': int(total_tokens),
            'lens_breakdown': lens_counts,
            'recent_sources': recent_sources
        }
    except Exception:
        return None


def clear_index():
    """Clear the entire vector index"""
    try:
        db = init_db()
        db.drop_table("roadmap_chunks")
        st.session_state.index_stats = None
        return True
    except Exception as e:
        st.error(f"Error clearing index: {e}")
        return False


def save_env_vars(anthropic_key: str, voyage_key: str):
    """Save API keys to .env file"""
    env_path = Path(".env")
    content = f"""# Anthropic Claude API Key (required)
# Get from: https://console.anthropic.com/
ANTHROPIC_API_KEY={anthropic_key}

# Voyage AI API Key (required)
# Get from: https://dash.voyageai.com/
VOYAGE_API_KEY={voyage_key}
"""
    env_path.write_text(content)
    # Reload environment
    os.environ['ANTHROPIC_API_KEY'] = anthropic_key
    os.environ['VOYAGE_API_KEY'] = voyage_key


def rebuild_context_graph():
    """Rebuild the context graph from all indexed chunks"""
    try:
        db = init_db()
        table = db.open_table("roadmap_chunks")

        # Get all chunks and embeddings from store
        all_data = table.to_pandas()
        all_chunks = []
        all_embeddings = []

        for _, row in all_data.iterrows():
            all_chunks.append({
                "id": row["id"],
                "content": row["content"],
                "lens": row["lens"],
                "source_file": row["source_file"],
                "chunk_index": row["chunk_index"],
                "token_count": row["token_count"],
            })
            all_embeddings.append(row["vector"])

        # Build graph
        graph = ContextGraph()
        graph.build_from_chunks(all_chunks, all_embeddings)
        graph.save()

        return graph.get_stats()
    except Exception as e:
        st.warning(f"⚠️ Could not build graph: {e}")
        return None


def get_all_materials() -> List[Dict]:
    """Get all materials files organized by lens"""
    materials = []
    for lens in VALID_LENSES:
        lens_dir = MATERIALS_DIR / lens
        if lens_dir.exists():
            for file_path in lens_dir.rglob("*"):
                if file_path.is_file() and not file_path.name.startswith('.'):
                    materials.append({
                        'file': file_path.name,
                        'path': str(file_path),
                        'lens': lens,
                        'size': file_path.stat().st_size,
                        'modified': datetime.fromtimestamp(file_path.stat().st_mtime),
                        'size_mb': f"{file_path.stat().st_size / 1024 / 1024:.2f} MB"
                    })
    return materials


def move_file_to_lens(file_path: str, new_lens: str) -> bool:
    """Move a file to a different lens folder"""
    try:
        source = Path(file_path)
        if not source.exists():
            return False

        # Create destination directory
        dest_dir = MATERIALS_DIR / new_lens
        dest_dir.mkdir(parents=True, exist_ok=True)

        # Move file
        dest = dest_dir / source.name
        source.rename(dest)
        return True
    except Exception as e:
        st.error(f"Error moving file: {e}")
        return False


def delete_material_file(file_path: str) -> bool:
    """Delete a material file"""
    try:
        Path(file_path).unlink()
        return True
    except Exception as e:
        st.error(f"Error deleting file: {e}")
        return False


def get_all_chunks() -> Optional[pd.DataFrame]:
    """Get all chunks from the index"""
    try:
        db = init_db()
        table = db.open_table("roadmap_chunks")
        chunks_df = table.to_pandas()

        if chunks_df.empty:
            return None

        # Add human-readable created_at
        chunks_df['created_at'] = pd.to_datetime(chunks_df['created_at'])
        chunks_df['created_at_str'] = chunks_df['created_at'].dt.strftime('%Y-%m-%d %H:%M')

        return chunks_df
    except Exception:
        return None


# ========== DASHBOARD COMPONENTS ==========

def render_quick_actions():
    """Render quick action buttons for common workflows."""

    st.subheader("⚡ Quick Actions")

    col1, col2, col3, col4 = st.columns(4)

    with col1:
        if st.button("🔄 Generate Roadmap", key="qa_gen_roadmap", use_container_width=True):
            with st.spinner("Generating roadmap..."):
                try:
                    result = generate_roadmap()
                    st.success("Roadmap generated!")
                    st.session_state.last_action = "roadmap_generated"
                except Exception as e:
                    st.error(f"Error: {e}")

    with col2:
        if st.button("🔄 Sync Graph", key="qa_sync_graph", use_container_width=True):
            with st.spinner("Syncing knowledge graph..."):
                try:
                    graph = sync_all_to_graph()
                    node_count = graph.graph.number_of_nodes() if graph.graph else 0
                    st.success(f"Graph synced! {node_count} nodes")
                except Exception as e:
                    st.error(f"Error: {e}")

    with col3:
        pending_count = len([q for q in load_questions() if q.get("status") == "pending"])
        if st.button(f"❓ Answer Questions ({pending_count})", key="qa_questions", use_container_width=True):
            st.session_state.current_page = "❓ Open Questions"
            st.rerun()

    with col4:
        if st.button("📄 Ingest Document", key="qa_ingest", use_container_width=True):
            st.session_state.current_page = "📁 Sources"
            st.session_state.show_ingest_form = True
            st.rerun()


def needs_graph_sync() -> bool:
    """Check if graph needs syncing based on file modification times."""
    try:
        graph_file = Path("data/unified_graph/graph.json")
        decisions_file = Path("data/questions/decisions.json")
        assessments_dir = Path("output/competitive")

        if not graph_file.exists():
            return True

        graph_mtime = graph_file.stat().st_mtime

        # Check if decisions are newer
        if decisions_file.exists() and decisions_file.stat().st_mtime > graph_mtime:
            return True

        # Check if any assessment is newer
        if assessments_dir.exists():
            for f in assessments_dir.glob("*.json"):
                if f.stat().st_mtime > graph_mtime:
                    return True

        return False
    except:
        return False


def render_attention_needed():
    """Render attention needed section on dashboard."""

    attention_items = []

    # Check for critical questions
    questions = load_questions()
    critical_pending = [q for q in questions if q.get("status") == "pending" and q.get("priority") == "critical"]
    if critical_pending:
        attention_items.append({
            "icon": "🔴",
            "message": f"{len(critical_pending)} critical question(s) pending",
            "action": "Go to Questions"
        })

    # Check for high priority pending questions
    high_pending = [q for q in questions if q.get("status") == "pending" and q.get("priority") == "high"]
    if high_pending:
        attention_items.append({
            "icon": "🟠",
            "message": f"{len(high_pending)} high priority question(s) pending",
            "action": "Go to Questions"
        })

    # Check for gaps without decisions
    try:
        graph = UnifiedContextGraph.load()
        if graph:
            unaddressed_gaps = [
                g for g in graph.node_indices.get("gap", {}).values()
                if not g.get("addressed_by_decision")
            ]
            critical_gaps = [g for g in unaddressed_gaps if g.get("severity") in ["critical", "significant"]]
            if critical_gaps:
                attention_items.append({
                    "icon": "⚠️",
                    "message": f"{len(critical_gaps)} significant gap(s) without decisions",
                    "action": "View Gaps"
                })
    except:
        pass

    # Check for graph sync needed
    if needs_graph_sync():
        attention_items.append({
            "icon": "🔄",
            "message": "Knowledge graph needs sync",
            "action": "Sync Now"
        })

    # Render
    if attention_items:
        st.subheader("⚡ Attention Needed")
        for item in attention_items:
            col1, col2 = st.columns([4, 1])
            col1.markdown(f"{item['icon']} {item['message']}")
            col2.button(item['action'], key=f"attn_{item['message'][:10]}")
    else:
        st.success("✅ No items need attention")


# ========== PAGE: DASHBOARD ==========

def page_dashboard():
    st.title("📊 Dashboard")

    # Quick Actions
    render_quick_actions()

    st.divider()

    st.markdown("Overview of your indexed materials and roadmap status")

    # Get or refresh stats
    if st.button("🔄 Refresh Stats"):
        st.session_state.index_stats = get_index_stats()

    stats = st.session_state.index_stats or get_index_stats()

    if not stats:
        st.warning("⚠️ No materials indexed yet. Go to **Ingest Materials** to get started.")
        return

    # Display key metrics
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Total Chunks", f"{stats['total_chunks']:,}")
    with col2:
        st.metric("Total Tokens", f"{stats['total_tokens']:,}")
    with col3:
        status = "✅ Ready to generate" if stats['total_chunks'] > 0 else "⚠️ No materials"
        st.metric("Status", status)

    st.divider()

    # Attention Needed
    render_attention_needed()

    st.divider()

    # Lens breakdown
    st.subheader("Breakdown by Lens")
    lens_df = pd.DataFrame([
        {"Lens": lens, "Chunks": count}
        for lens, count in stats['lens_breakdown'].items()
    ])
    st.bar_chart(lens_df.set_index("Lens"))

    # Recent sources
    st.subheader("Recent Ingested Sources")
    if stats['recent_sources']:
        recent_df = pd.DataFrame(stats['recent_sources'])
        recent_df['created_at'] = pd.to_datetime(recent_df['created_at'])
        recent_df['created_at'] = recent_df['created_at'].dt.strftime('%Y-%m-%d %H:%M')
        st.dataframe(recent_df, use_container_width=True)

    # Context graph stats
    st.subheader("Context Graph")
    try:
        graph = ContextGraph().load()
        graph_stats = graph.get_stats()

        col1, col2, col3 = st.columns(3)
        col1.metric("Graph Nodes", graph_stats["nodes"])
        col2.metric("Graph Edges", graph_stats["edges"])
        col3.metric("Components", graph_stats["components"])

        if graph_stats["edge_types"]:
            with st.expander("Edge Type Breakdown"):
                for edge_type, count in graph_stats["edge_types"].items():
                    st.write(f"- **{edge_type}**: {count}")
    except Exception:
        st.info("No context graph built yet. Ingest materials to build the graph.")

    # Unified Context Graph stats
    st.subheader("🕸️ Unified Knowledge Graph")
    st.caption("Integrates decisions, assessments, questions, and roadmap with authority hierarchy")

    try:
        unified_graph = UnifiedContextGraph.load()

        if unified_graph.graph.number_of_nodes() == 0:
            st.info("Graph is empty. Click 'Sync Graph' to build the unified knowledge graph.")

            if st.button("🔄 Sync Graph"):
                with st.spinner("Syncing all data sources..."):
                    unified_graph = sync_all_to_graph()
                st.success("Graph synced!")
                st.rerun()
        else:
            # Show graph metrics
            col1, col2, col3, col4 = st.columns(4)
            col1.metric("Total Nodes", unified_graph.graph.number_of_nodes())
            col2.metric("Total Edges", unified_graph.graph.number_of_edges())

            active_decisions = len([
                d for d in unified_graph.node_indices["decision"].values()
                if d.get("status") == "active"
            ])
            col3.metric("Active Decisions", active_decisions)

            pending_questions = len([
                q for q in unified_graph.node_indices["question"].values()
                if q.get("status") == "pending"
            ])
            col4.metric("Open Questions", pending_questions)

            # Authority breakdown
            with st.expander("Knowledge by Authority Level"):
                authority_data = {}
                for node_type, level in AUTHORITY_LEVELS.items():
                    count = 0
                    if node_type == "answered_question":
                        count = len([q for q in unified_graph.node_indices["question"].values() if q.get("status") == "answered"])
                    elif node_type == "pending_question":
                        count = len([q for q in unified_graph.node_indices["question"].values() if q.get("status") == "pending"])
                    else:
                        count = len(unified_graph.node_indices.get(node_type.replace("_question", "question"), {}))

                    if count > 0:
                        authority_data[f"L{level}: {node_type.replace('_', ' ').title()}"] = count

                if authority_data:
                    st.bar_chart(authority_data)

            if st.button("🔄 Re-sync Graph"):
                with st.spinner("Re-syncing..."):
                    unified_graph = sync_all_to_graph()
                st.success("Graph re-synced!")
                st.rerun()

    except Exception as e:
        st.warning(f"Could not load unified graph: {e}")
        if st.button("🔄 Build Graph"):
            with st.spinner("Building unified knowledge graph..."):
                unified_graph = sync_all_to_graph()
            st.success("Graph built!")
            st.rerun()

    # Clear index button
    st.divider()
    if st.button("🗑️ Clear Index", type="secondary"):
        if st.checkbox("⚠️ Are you sure? This cannot be undone."):
            if clear_index():
                st.success("Index cleared successfully!")
                st.rerun()


# ========== QUESTIONS PAGE COMPONENTS ==========

def get_decision_overrides(decision_id: str) -> list:
    """Get chunks that a decision overrides from the graph."""

    try:
        graph = UnifiedContextGraph.load()
        if not graph or not graph.graph:
            return []

        overrides = []

        # Find OVERRIDES edges from this decision
        if decision_id in graph.graph:
            for neighbor in graph.graph.successors(decision_id):
                edge_data = graph.graph.edges.get((decision_id, neighbor), {})
                if edge_data.get("edge_type") == "OVERRIDES":
                    # Get chunk data
                    chunk = graph.node_indices.get("chunk", {}).get(neighbor)
                    if chunk:
                        chunk_dict = chunk if isinstance(chunk, dict) else (chunk.__dict__ if hasattr(chunk, '__dict__') else {})
                        overrides.append(chunk_dict)

        return overrides
    except Exception as e:
        return []


def save_decision_update(decision: dict):
    """Update a single decision in the decisions file."""
    decisions = load_decisions()
    for i, d in enumerate(decisions):
        if d["id"] == decision["id"]:
            decisions[i] = decision
            break
    save_decisions(decisions)


def extract_key_terms_simple(text: str) -> list:
    """Simple keyword extraction without external dependencies."""
    import re
    from collections import Counter

    # Common stop words to filter out
    stop_words = {
        "the", "a", "an", "is", "are", "was", "were", "be", "been", "being",
        "have", "has", "had", "do", "does", "did", "will", "would", "could",
        "should", "may", "might", "can", "this", "that", "these", "those",
        "i", "you", "he", "she", "it", "we", "they", "what", "which", "who",
        "when", "where", "why", "how", "all", "each", "every", "both", "few",
        "more", "most", "other", "some", "such", "no", "not", "only", "own",
        "same", "so", "than", "too", "very", "just", "also", "now", "here",
        "there", "then", "once", "and", "or", "but", "if", "because", "as",
        "until", "while", "of", "at", "by", "for", "with", "about", "against",
        "between", "into", "through", "during", "before", "after", "above",
        "below", "to", "from", "up", "down", "in", "out", "on", "off", "over",
        "under", "again", "further", "our", "your", "their", "its"
    }

    # Tokenize (simple split, remove punctuation)
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())

    # Filter stop words and short words
    keywords = [w for w in words if w not in stop_words and len(w) > 2]

    # Count frequency and return top terms
    counts = Counter(keywords)

    # Return top 8 most frequent terms
    return [word for word, count in counts.most_common(8)]


def get_embedding(text: str) -> list:
    """Get embedding for a single text using Voyage AI."""
    try:
        embeddings = generate_embeddings([text])
        return embeddings[0] if embeddings else []
    except:
        return []


def get_lancedb_table():
    """Get LanceDB table for chunks."""
    try:
        db = init_db()
        return db.open_table("roadmap_chunks")
    except:
        return None


def search_chunks_semantic(query: str, limit: int = 10) -> list:
    """Semantic search on chunks using embeddings."""
    try:
        # Get embedding for query
        query_embedding = get_embedding(query)
        if not query_embedding:
            return []

        # Search LanceDB
        table = get_lancedb_table()
        if table is None:
            return []

        results = table.search(query_embedding).limit(limit).to_list()
        return results
    except Exception as e:
        return []


def search_chunks_keyword(keywords: list, limit: int = 10) -> list:
    """Keyword-based search on chunks."""
    try:
        table = get_lancedb_table()
        if table is None:
            return []

        # Get all chunks and filter
        df = table.to_pandas()

        results = []
        for _, row in df.iterrows():
            content = str(row.get("content", "")).lower()
            matches = sum(1 for kw in keywords if kw.lower() in content)

            if matches > 0:
                results.append({
                    "id": row.get("id", ""),
                    "source_name": row.get("source_file", "Unknown").split("/")[-1] if "/" in str(row.get("source_file", "")) else row.get("source_file", "Unknown"),
                    "lens": row.get("lens", ""),
                    "content": row.get("content", ""),
                    "similarity": matches / len(keywords) if keywords else 0,
                    "matched_count": matches
                })

        # Sort by match count
        results.sort(key=lambda x: x["matched_count"], reverse=True)

        return results[:limit]
    except Exception as e:
        return []


def find_sources_for_question(question: dict, max_sources: int = 5) -> list:
    """Dynamically find source references for a question by searching graph and chunks."""

    question_text = question.get("question", "")
    context = question.get("context", "")
    query = f"{question_text} {context}".strip()

    if not query:
        return []

    sources = []
    seen_ids = set()

    # === METHOD 1: Semantic search on chunks (LanceDB) ===
    try:
        chunk_results = search_chunks_semantic(query, limit=10)

        for chunk in chunk_results:
            chunk_id = chunk.get("id", "")
            if chunk_id in seen_ids or not chunk_id:
                continue
            seen_ids.add(chunk_id)

            source_file = chunk.get("source_file", "Unknown")
            source_name = source_file.split("/")[-1] if "/" in source_file else source_file

            sources.append({
                "type": "chunk",
                "id": chunk_id,
                "source_name": source_name,
                "source_path": chunk.get("source_file", ""),
                "lens": chunk.get("lens", "unknown"),
                "content": chunk.get("content", "")[:300],
                "similarity": chunk.get("_distance", 0),
                "search_method": "semantic"
            })
    except Exception as e:
        pass

    # === METHOD 2: Keyword extraction and search ===
    try:
        keywords = extract_key_terms_simple(question_text)

        if keywords:
            keyword_results = search_chunks_keyword(keywords, limit=10)

            for chunk in keyword_results:
                chunk_id = chunk.get("id", "")
                if chunk_id in seen_ids or not chunk_id:
                    continue
                seen_ids.add(chunk_id)

                sources.append({
                    "type": "chunk",
                    "id": chunk_id,
                    "source_name": chunk.get("source_name", "Unknown"),
                    "source_path": chunk.get("source_file", ""),
                    "lens": chunk.get("lens", "unknown"),
                    "content": chunk.get("content", "")[:300],
                    "similarity": chunk.get("similarity", 0),
                    "search_method": "keyword",
                    "matched_terms": [kw for kw in keywords if kw.lower() in chunk.get("content", "").lower()]
                })
    except Exception as e:
        pass

    # === METHOD 3: Graph traversal ===
    try:
        graph = UnifiedContextGraph.load()

        if graph:
            keywords = extract_key_terms_simple(question_text)

            # Search assessments
            for assess_id, assessment in graph.node_indices.get("assessment", {}).items():
                if assess_id in seen_ids:
                    continue

                summary = ""
                if isinstance(assessment, dict):
                    summary = assessment.get("summary", "")
                else:
                    summary = getattr(assessment, 'summary', "")

                if any(kw.lower() in summary.lower() for kw in keywords):
                    seen_ids.add(assess_id)

                    assess_type = ""
                    if isinstance(assessment, dict):
                        assess_type = assessment.get("assessment_type", "Unknown")
                    else:
                        assess_type = getattr(assessment, 'assessment_type', "Unknown")

                    sources.append({
                        "type": "assessment",
                        "id": assess_id,
                        "source_name": f"{assess_type.title()} Assessment",
                        "lens": "assessment",
                        "content": summary[:300],
                        "similarity": 0.7,
                        "search_method": "graph"
                    })

            # Search gaps
            for gap_id, gap in graph.node_indices.get("gap", {}).items():
                if gap_id in seen_ids:
                    continue

                description = ""
                if isinstance(gap, dict):
                    description = gap.get("description", "")
                else:
                    description = getattr(gap, 'description', "")

                if any(kw.lower() in description.lower() for kw in keywords):
                    seen_ids.add(gap_id)
                    sources.append({
                        "type": "gap",
                        "id": gap_id,
                        "source_name": f"Gap: {description[:50]}",
                        "lens": "gap",
                        "content": description[:300],
                        "similarity": 0.6,
                        "search_method": "graph"
                    })
    except Exception as e:
        pass

    # === RANK AND RETURN TOP SOURCES ===
    sources.sort(key=lambda x: x.get("similarity", 0), reverse=True)
    return sources[:max_sources]


def render_question_source_references(question: dict):
    """Render source references for a question. Dynamically finds sources if not cached."""

    # Check for cached sources
    cache_key = f"sources_{question['id']}"

    if cache_key in st.session_state:
        sources = st.session_state[cache_key]
    else:
        # Find sources dynamically
        with st.spinner("Finding source references..."):
            sources = find_sources_for_question(question, max_sources=5)
            st.session_state[cache_key] = sources

    if not sources:
        st.caption("No relevant sources found")

        # Offer to refresh
        if st.button("🔄 Search Again", key=f"refresh_sources_{question['id']}"):
            if cache_key in st.session_state:
                del st.session_state[cache_key]
            st.rerun()
        return

    st.markdown(f"**📚 Source References** ({len(sources)} found)")

    for i, source in enumerate(sources):
        source_type = source.get("type", "chunk")
        source_id = source.get("id", "unknown")
        source_name = source.get("source_name", "Unknown Source")
        lens = source.get("lens", "")
        content = source.get("content", "")
        similarity = source.get("similarity", 0)
        search_method = source.get("search_method", "unknown")
        matched_terms = source.get("matched_terms", [])

        # Type icon
        type_icons = {
            "chunk": "📄",
            "assessment": "🔬",
            "roadmap_item": "🗺️",
            "gap": "⚠️",
            "decision": "✅"
        }
        icon = type_icons.get(source_type, "📎")

        # Method badge
        method_badges = {
            "semantic": "🎯",
            "keyword": "🔤",
            "graph": "🕸️"
        }
        method_badge = method_badges.get(search_method, "")

        with st.container(border=True):
            # Header row
            col1, col2 = st.columns([4, 1])

            with col1:
                st.markdown(f"{icon} **{source_name}**")

                meta_parts = []
                if lens and lens not in ["assessment", "gap", "roadmap"]:
                    meta_parts.append(f"Lens: {lens}")
                meta_parts.append(f"ID: {source_id[:20]}...")
                meta_parts.append(f"{method_badge} {search_method}")

                st.caption(" | ".join(meta_parts))

            with col2:
                # Relevance score
                score_pct = min(similarity * 100, 100)
                if score_pct >= 70:
                    st.success(f"{score_pct:.0f}%")
                elif score_pct >= 40:
                    st.warning(f"{score_pct:.0f}%")
                else:
                    st.caption(f"{score_pct:.0f}%")

            # Content excerpt
            if content:
                display_content = content[:250]
                if len(content) > 250:
                    display_content += "..."

                st.markdown(f"*\"{display_content}\"*")

            # Matched keywords (for keyword search)
            if matched_terms:
                st.caption(f"Matched: {', '.join(matched_terms)}")

            # === VIEW ORIGINAL DOCUMENT ===
            if source_type == "chunk" and source.get("source_path"):
                with st.expander("📄 View Original Document"):
                    render_original_document_viewer(source, content, unique_key=f"{question['id']}_{i}")

            elif source_type == "assessment":
                with st.expander("🔬 View Full Assessment"):
                    render_assessment_detail(source_id)

            elif source_type == "roadmap_item":
                with st.expander("🗺️ View Roadmap Item"):
                    render_roadmap_item_detail(source_id)

    # Refresh button
    col1, col2 = st.columns([3, 1])
    with col2:
        if st.button("🔄 Refresh", key=f"refresh_sources_{question['id']}"):
            if cache_key in st.session_state:
                del st.session_state[cache_key]
            st.rerun()


def invalidate_source_cache():
    """Clear all cached source references."""
    keys_to_delete = [k for k in st.session_state.keys() if k.startswith("sources_")]
    for key in keys_to_delete:
        del st.session_state[key]


def get_original_document(source_path: str) -> dict | None:
    """
    Load original document from source path.

    Returns dict with content, metadata, or None if not found.
    """
    from pathlib import Path
    import os

    if not source_path:
        return None

    path = Path(source_path)

    # Handle relative paths
    if not path.is_absolute():
        # Try common base paths
        for base in [".", "materials", os.path.expanduser("~/roadmap-synth")]:
            full_path = Path(base) / path
            if full_path.exists():
                path = full_path
                break

    if not path.exists():
        return None

    try:
        # Get file metadata
        stat = path.stat()
        modified = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")
        size = stat.st_size

        # Read content based on file type
        suffix = path.suffix.lower()

        if suffix in [".txt", ".md", ".json", ".csv", ".yaml", ".yml"]:
            # Plain text files
            content = path.read_text(encoding="utf-8", errors="replace")
            content_type = "text"

        elif suffix in [".pdf"]:
            # PDF - extract text if possible
            try:
                import pypdf
                reader = pypdf.PdfReader(str(path))
                content = "\n\n".join(page.extract_text() for page in reader.pages)
                content_type = "pdf_extracted"
            except:
                content = f"[PDF file - {size} bytes - text extraction not available]"
                content_type = "binary"

        elif suffix in [".docx"]:
            # Word doc - extract text if possible
            try:
                import docx
                doc = docx.Document(str(path))
                content = "\n\n".join(para.text for para in doc.paragraphs)
                content_type = "docx_extracted"
            except:
                content = f"[Word document - {size} bytes - text extraction not available]"
                content_type = "binary"

        elif suffix in [".pptx"]:
            # PowerPoint - extract text if possible
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                content = "\n\n".join(texts)
                content_type = "pptx_extracted"
            except:
                content = f"[PowerPoint - {size} bytes - text extraction not available]"
                content_type = "binary"

        else:
            # Unknown type - try to read as text
            try:
                content = path.read_text(encoding="utf-8", errors="replace")
                content_type = "text"
            except:
                content = f"[Binary file - {size} bytes]"
                content_type = "binary"

        return {
            "path": str(path),
            "name": path.name,
            "size": size,
            "modified": modified,
            "content": content,
            "content_type": content_type,
            "suffix": suffix
        }

    except Exception as e:
        print(f"Error reading document: {e}")
        return None


def find_chunk_in_document(document_content: str, chunk_content: str) -> dict | None:
    """
    Find the approximate location of a chunk within the original document.

    Returns dict with line numbers and context, or None if not found.
    """
    import re

    if not document_content or not chunk_content:
        return None

    # Clean up for matching
    doc_lower = document_content.lower()

    # Try to find the chunk content (first 100 chars for matching)
    search_text = chunk_content[:100].lower().strip()

    # Remove extra whitespace for fuzzy matching
    search_text_normalized = re.sub(r'\s+', ' ', search_text)
    doc_normalized = re.sub(r'\s+', ' ', doc_lower)

    pos = doc_normalized.find(search_text_normalized[:50])

    if pos == -1:
        # Try with even shorter match
        pos = doc_normalized.find(search_text_normalized[:30])

    if pos == -1:
        return None

    # Calculate approximate line number
    text_before = document_content[:pos]
    start_line = text_before.count('\n') + 1

    # Estimate end line
    chunk_lines = chunk_content.count('\n')
    end_line = start_line + chunk_lines

    return {
        "start_line": start_line,
        "end_line": end_line,
        "char_position": pos
    }


def render_original_document_viewer(source: dict, chunk_content: str, unique_key: str = ""):
    """
    Render a viewer for the original document with chunk highlighting.
    """
    source_path = source.get("source_path", "")

    if not source_path:
        st.caption("Original document path not available")
        return

    # Load document
    doc = get_original_document(source_path)

    if not doc:
        st.warning(f"Could not load original document: {source_path}")
        return

    # Find chunk location
    location = find_chunk_in_document(doc["content"], chunk_content)

    # Document header
    st.markdown(f"**📄 Original Document:** `{doc['name']}`")

    col1, col2, col3 = st.columns(3)
    col1.caption(f"Size: {doc['size']:,} bytes")
    col2.caption(f"Modified: {doc['modified']}")
    col3.caption(f"Type: {doc['content_type']}")

    if location:
        st.caption(f"📍 Chunk location: Lines {location['start_line']}-{location['end_line']} (approx)")

    # Content display
    content = doc["content"]

    # Limit display size for very large files
    MAX_DISPLAY_CHARS = 50000
    truncated = False

    if len(content) > MAX_DISPLAY_CHARS:
        # If we know chunk location, center around it
        if location:
            start_char = max(0, location["char_position"] - 5000)
            end_char = min(len(content), location["char_position"] + 10000)
            content = f"... [truncated - showing around chunk location] ...\n\n{content[start_char:end_char]}\n\n... [truncated] ..."
        else:
            content = content[:MAX_DISPLAY_CHARS] + f"\n\n... [truncated - {len(doc['content']) - MAX_DISPLAY_CHARS:,} more characters]"
        truncated = True

    # Display content in scrollable container
    st.text_area(
        "Document Content",
        value=content,
        height=400,
        disabled=True,
        key=f"doc_content_{unique_key}" if unique_key else f"doc_content_{source.get('id', 'unknown')}"
    )

    if truncated:
        st.caption("Document truncated for display. Download for full content.")

    # Download button
    st.download_button(
        "📥 Download Full Document",
        doc["content"],
        file_name=doc["name"],
        mime="text/plain",
        key=f"download_{unique_key}" if unique_key else f"download_{source.get('id', 'unknown')}"
    )


def render_assessment_detail(assessment_id: str):
    """Render full assessment detail."""
    try:
        graph = load_unified_graph()
        if not graph:
            st.caption("Graph not available")
            return

        assessment = graph.node_indices.get("assessment", {}).get(assessment_id)

        if not assessment:
            st.caption(f"Assessment {assessment_id} not found")
            return

        data = assessment.__dict__ if hasattr(assessment, '__dict__') else assessment

        st.markdown(f"**Type:** {data.get('assessment_type', 'Unknown')}")
        st.markdown(f"**Created:** {data.get('created_at', 'Unknown')[:10]}")
        st.markdown(f"**Confidence:** {data.get('confidence', 'Unknown')}")

        st.divider()

        st.markdown("**Summary:**")
        st.write(data.get("summary", "No summary available"))

        # Show full assessment data
        if data.get("assessment_data"):
            with st.expander("Raw Assessment Data"):
                st.json(data["assessment_data"])

    except Exception as e:
        st.error(f"Error loading assessment: {e}")


def render_roadmap_item_detail(item_id: str):
    """Render roadmap item detail."""
    try:
        graph = load_unified_graph()
        if not graph:
            st.caption("Graph not available")
            return

        item = graph.node_indices.get("roadmap_item", {}).get(item_id)

        if not item:
            st.caption(f"Roadmap item {item_id} not found")
            return

        data = item.__dict__ if hasattr(item, '__dict__') else item

        st.markdown(f"### {data.get('name', 'Unknown')}")
        st.markdown(f"**Horizon:** {data.get('horizon', 'Unknown')}")
        st.markdown(f"**Theme:** {data.get('theme', 'Unknown')}")
        st.markdown(f"**Owner:** {data.get('owner', 'Unknown')}")

        st.divider()

        st.markdown("**Description:**")
        st.write(data.get("description", "No description available"))

        # Show gaps if any
        gaps = data.get("has_gaps", [])
        if gaps:
            st.markdown(f"**Gaps:** {len(gaps)}")
            for gap_id in gaps[:5]:
                st.caption(f"- {gap_id}")

        # Show questions if any
        questions = data.get("has_questions", [])
        if questions:
            st.markdown(f"**Open Questions:** {len(questions)}")
            for q_id in questions[:5]:
                st.caption(f"- {q_id}")

    except Exception as e:
        st.error(f"Error loading roadmap item: {e}")


def save_question_validation(question_id: str, is_accurate: bool, validated_by: str, feedback_note: str = ""):
    """Save validation feedback for a question."""

    questions = load_questions()

    for q in questions:
        if q["id"] == question_id:
            q["validation"] = {
                "validated": True,
                "is_accurate": is_accurate,
                "validated_by": validated_by,
                "validated_at": datetime.now().isoformat(),
                "feedback_note": feedback_note
            }
            break

    save_questions(questions)


def render_question_validation(question: dict):
    """Render validation feedback UI for a question."""

    validation = question.get("validation")

    st.markdown("**✅ Validate Question**")

    if validation and validation.get("validated"):
        # Already validated - show result
        is_accurate = validation.get("is_accurate")
        validated_by = validation.get("validated_by", "Unknown")
        validated_at = validation.get("validated_at", "")[:10]

        if is_accurate:
            st.success(f"👍 Validated as accurate by {validated_by} on {validated_at}")
        else:
            st.error(f"👎 Validated as inaccurate by {validated_by} on {validated_at}")

        if validation.get("feedback_note"):
            st.caption(f"Note: {validation['feedback_note']}")

        # Option to re-validate
        if st.button("Re-validate", key=f"reval_{question['id']}"):
            st.session_state.revalidating_question = question["id"]
            st.rerun()

    else:
        # Not yet validated - show validation UI
        st.caption("Review the source references above, then indicate if this question is well-grounded:")

        col1, col2, col3 = st.columns([1, 1, 2])

        with col1:
            thumbs_up = st.button("👍 Accurate", key=f"thumbs_up_{question['id']}", use_container_width=True)

        with col2:
            thumbs_down = st.button("👎 Inaccurate", key=f"thumbs_down_{question['id']}", use_container_width=True)

        with col3:
            feedback_note = st.text_input(
                "Feedback (optional)",
                key=f"feedback_{question['id']}",
                placeholder="Why accurate/inaccurate?"
            )

        validated_by = st.text_input(
            "Your name",
            key=f"validator_{question['id']}",
            value=st.session_state.get("last_validator", "")
        )

        if thumbs_up or thumbs_down:
            if not validated_by:
                st.warning("Please enter your name")
            else:
                # Save validation
                save_question_validation(
                    question_id=question["id"],
                    is_accurate=thumbs_up,
                    validated_by=validated_by,
                    feedback_note=feedback_note
                )
                st.session_state.last_validator = validated_by
                st.success("Validation saved!")
                st.rerun()


def render_validation_stats():
    """Render validation statistics."""

    questions = load_questions()

    # Filter out None values and ensure dict structure
    questions = [q for q in questions if q is not None and isinstance(q, dict)]

    total = len(questions)
    validated = [q for q in questions if isinstance(q.get("validation"), dict) and q.get("validation", {}).get("validated")]
    accurate = [q for q in validated if q.get("validation", {}).get("is_accurate")]
    inaccurate = [q for q in validated if not q.get("validation", {}).get("is_accurate")]

    st.subheader("📊 Question Validation Stats")

    col1, col2, col3, col4 = st.columns(4)

    col1.metric("Total Questions", total)
    col2.metric("Validated", len(validated))
    col3.metric("👍 Accurate", len(accurate))
    col4.metric("👎 Inaccurate", len(inaccurate))

    if validated:
        accuracy_rate = len(accurate) / len(validated) * 100
        st.progress(accuracy_rate / 100)
        st.caption(f"Accuracy Rate: {accuracy_rate:.1f}%")

    # Recent inaccurate (for review)
    if inaccurate:
        with st.expander(f"⚠️ Review Inaccurate Questions ({len(inaccurate)})"):
            for q in inaccurate[:5]:
                st.markdown(f"- **{q['question'][:60]}...**")
                st.caption(f"  Feedback: {q['validation'].get('feedback_note', 'No feedback')}")


# ========== HOLISTIC QUESTION GENERATION ==========

def cosine_similarity(vec1: list, vec2: list) -> float:
    """Calculate cosine similarity between two vectors."""
    import math

    if not vec1 or not vec2 or len(vec1) != len(vec2):
        return 0.0

    dot_product = sum(a * b for a, b in zip(vec1, vec2))
    magnitude1 = math.sqrt(sum(a * a for a in vec1))
    magnitude2 = math.sqrt(sum(b * b for b in vec2))

    if magnitude1 == 0 or magnitude2 == 0:
        return 0.0

    return dot_product / (magnitude1 * magnitude2)


def gather_generation_context() -> dict:
    """Gather all context needed for question generation."""

    context = {
        # Questions
        "existing_questions": load_questions(),
        "pending_questions": [q for q in load_questions() if q.get("status") == "pending"],
        "answered_questions": [q for q in load_questions() if q.get("status") == "answered"],

        # Decisions
        "active_decisions": [d for d in load_decisions() if d.get("status") == "active"],

        # Graph
        "graph": UnifiedContextGraph.load(),

        # Assessments
        "arch_assessments": load_alignment_analysis() or [],
        "competitive_assessments": load_analyst_assessments() or [],

        # Roadmap
        "roadmap": None,  # Will be loaded from graph
        "roadmap_items": [],

        # Source stats
        "chunks_by_lens": get_chunks_by_lens_count(),
    }

    # Extract roadmap items from graph
    graph = context["graph"]
    if graph:
        roadmap_items_nodes = graph.node_indices.get("roadmap_item", {})
        context["roadmap_items"] = []
        for item_id, item_node in roadmap_items_nodes.items():
            item_data = item_node.__dict__ if hasattr(item_node, '__dict__') else item_node
            context["roadmap_items"].append({
                "id": item_id,
                "name": item_data.get("name", "Unknown"),
                "horizon": item_data.get("horizon", "future"),
                "owner": item_data.get("owner", ""),
                "dependencies": item_data.get("dependencies", [])
            })

    return context


def find_contradictions(graph) -> list[dict]:
    """
    Find contradictory statements in the graph.

    Looks for chunks with topic overlap but different lenses.
    """

    contradictions = []

    if not graph:
        return contradictions

    chunks = list(graph.node_indices.get("chunk", {}).values())

    # Compare your-voice chunks against team chunks
    your_voice_chunks = [c for c in chunks if getattr(c, 'lens', c.get('lens')) == 'your-voice']
    team_chunks = [c for c in chunks if getattr(c, 'lens', c.get('lens')) in ['team-structured', 'team-conversational', 'engineering']]

    # Look for topic overlap with different statements
    contradiction_keywords = [
        ("q1", "q2", "q3", "q4"),  # Timeline conflicts
        ("will", "won't", "can", "cannot", "can't"),  # Capability conflicts
        ("priority", "deprioritize"),  # Priority conflicts
        ("before", "after"),  # Sequence conflicts
    ]

    for yv_chunk in your_voice_chunks[:20]:  # Limit for performance
        yv_content = getattr(yv_chunk, 'content', yv_chunk.get('content', '')).lower()
        yv_id = getattr(yv_chunk, 'id', yv_chunk.get('id', ''))
        yv_source = getattr(yv_chunk, 'source_name', yv_chunk.get('source_name', ''))

        for team_chunk in team_chunks[:50]:
            team_content = getattr(team_chunk, 'content', team_chunk.get('content', '')).lower()
            team_id = getattr(team_chunk, 'id', team_chunk.get('id', ''))
            team_source = getattr(team_chunk, 'source_name', team_chunk.get('source_name', ''))
            team_lens = getattr(team_chunk, 'lens', team_chunk.get('lens', ''))

            # Check for topic overlap (shared important terms)
            yv_terms = set(extract_key_terms_simple(yv_content))
            team_terms = set(extract_key_terms_simple(team_content))

            overlap = yv_terms & team_terms

            if len(overlap) >= 3:  # Significant topic overlap
                # Check for potential contradiction signals
                for keyword_group in contradiction_keywords:
                    yv_has = any(kw in yv_content for kw in keyword_group)
                    team_has = any(kw in team_content for kw in keyword_group)

                    if yv_has and team_has:
                        # Potential contradiction found
                        topic = ", ".join(list(overlap)[:3])

                        contradictions.append({
                            "topic": topic,
                            "chunk_a_id": yv_id,
                            "source_a": yv_source,
                            "lens_a": "your-voice",
                            "statement_a": getattr(yv_chunk, 'content', yv_chunk.get('content', ''))[:150],
                            "chunk_b_id": team_id,
                            "source_b": team_source,
                            "lens_b": team_lens,
                            "statement_b": getattr(team_chunk, 'content', team_chunk.get('content', ''))[:150],
                        })
                        break  # One contradiction per chunk pair

    # Deduplicate similar contradictions
    seen_topics = set()
    unique_contradictions = []
    for c in contradictions:
        topic_key = c["topic"]
        if topic_key not in seen_topics:
            seen_topics.add(topic_key)
            unique_contradictions.append(c)

    return unique_contradictions[:5]  # Limit to top 5


def derive_questions_from_graph(context: dict) -> list[dict]:
    """
    Automatically derive questions from graph analysis.
    These are flagged as 'derived' type.
    """

    derived_questions = []
    graph = context["graph"]
    roadmap_items = context["roadmap_items"]
    decisions = context["active_decisions"]

    if not graph:
        return derived_questions

    # === PATTERN 1: Contradictions ===
    contradictions = find_contradictions(graph)
    for contradiction in contradictions:
        derived_questions.append({
            "question": f"Conflicting information about {contradiction['topic']}: "
                       f"'{contradiction['statement_a'][:40]}...' vs "
                       f"'{contradiction['statement_b'][:40]}...'. Which is correct?",
            "audience": "leadership",
            "category": "alignment",
            "priority": "high",
            "context": f"Found in {contradiction['source_a']} ({contradiction['lens_a']}) "
                      f"and {contradiction['source_b']} ({contradiction['lens_b']})",
            "generation": {
                "type": "derived",
                "source": "contradiction",
                "generated_at": datetime.now().isoformat(),
            },
            "derivation": {
                "pattern": "contradiction",
                "evidence": [
                    {
                        "source_id": contradiction["chunk_a_id"],
                        "source_name": contradiction["source_a"],
                        "lens": contradiction["lens_a"],
                        "content": contradiction["statement_a"],
                    },
                    {
                        "source_id": contradiction["chunk_b_id"],
                        "source_name": contradiction["source_b"],
                        "lens": contradiction["lens_b"],
                        "content": contradiction["statement_b"],
                    }
                ]
            }
        })

    # === PATTERN 2: Missing Engineering Coverage ===
    for item in roadmap_items[:10]:  # Limit for performance
        # Get supporting chunks for this item
        item_name_lower = item['name'].lower()
        supporting_chunks = []

        chunks = list(graph.node_indices.get("chunk", {}).values())
        for chunk in chunks:
            chunk_content = getattr(chunk, 'content', chunk.get('content', '')).lower()
            if any(word in chunk_content for word in item_name_lower.split()[:3]):
                supporting_chunks.append(chunk)

        has_engineering = any(getattr(c, 'lens', c.get('lens')) == 'engineering' for c in supporting_chunks)

        if not has_engineering and item.get("horizon") in ["now", "next"]:
            derived_questions.append({
                "question": f"Has engineering validated the feasibility of '{item['name']}'?",
                "audience": "engineering",
                "category": "feasibility",
                "priority": "high" if item.get("horizon") == "now" else "medium",
                "context": f"Roadmap item '{item['name']}' in {item.get('horizon')} horizon "
                          f"has no engineering source documents",
                "generation": {
                    "type": "derived",
                    "source": "missing_coverage",
                    "generated_at": datetime.now().isoformat(),
                },
                "derivation": {
                    "pattern": "missing_coverage",
                    "evidence": [{
                        "roadmap_item": item["name"],
                        "horizon": item.get("horizon"),
                        "lenses_present": list(set(getattr(c, 'lens', c.get('lens')) for c in supporting_chunks)) if supporting_chunks else [],
                        "missing_lens": "engineering"
                    }]
                }
            })

    # === PATTERN 3: Dependency Conflicts ===
    horizon_order = {"now": 0, "next": 1, "later": 2, "future": 3}

    for item in roadmap_items[:10]:
        dependencies = item.get("dependencies", [])
        item_horizon = item.get("horizon", "future")

        for dep_name in dependencies:
            dep_item = next((i for i in roadmap_items if i["name"] == dep_name), None)
            if dep_item:
                dep_horizon = dep_item.get("horizon", "future")

                # Check if dependency is in a later horizon
                if horizon_order.get(dep_horizon, 3) > horizon_order.get(item_horizon, 3):
                    derived_questions.append({
                        "question": f"'{item['name']}' is in {item_horizon} but depends on "
                                   f"'{dep_name}' which is in {dep_horizon}. How will this work?",
                        "audience": "product",
                        "category": "dependency",
                        "priority": "high",
                        "context": f"Dependency timing conflict detected",
                        "generation": {
                            "type": "derived",
                            "source": "dependency_conflict",
                            "generated_at": datetime.now().isoformat(),
                        },
                        "derivation": {
                            "pattern": "dependency_conflict",
                            "evidence": [{
                                "item": item["name"],
                                "item_horizon": item_horizon,
                                "dependency": dep_name,
                                "dependency_horizon": dep_horizon
                            }]
                        }
                    })

    # === PATTERN 4: Unaddressed Gaps ===
    gaps = list(graph.node_indices.get("gap", {}).values())[:10]  # Limit for performance
    for gap in gaps:
        gap_data = gap.__dict__ if hasattr(gap, '__dict__') else gap

        # Check if gap has been addressed by a decision
        addressed_by = gap_data.get("addressed_by_decision")
        if not addressed_by:
            severity = gap_data.get("severity", "medium")
            if severity in ["critical", "significant"]:
                derived_questions.append({
                    "question": f"How should we address the gap: '{gap_data.get('description', 'Unknown')[:80]}'?",
                    "audience": "leadership",
                    "category": "direction",
                    "priority": "critical" if severity == "critical" else "high",
                    "context": f"Gap identified by assessment, severity: {severity}",
                    "generation": {
                        "type": "derived",
                        "source": "unaddressed_gap",
                        "generated_at": datetime.now().isoformat(),
                    },
                    "derivation": {
                        "pattern": "unaddressed_gap",
                        "evidence": [{
                            "gap_id": gap_data.get("id"),
                            "description": gap_data.get("description"),
                            "severity": severity
                        }]
                    }
                })

    # === PATTERN 5: Missing Owner ===
    for item in roadmap_items[:10]:
        if not item.get("owner") and item.get("horizon") in ["now", "next"]:
            derived_questions.append({
                "question": f"Who owns delivery of '{item['name']}'?",
                "audience": "leadership",
                "category": "ownership",
                "priority": "high" if item.get("horizon") == "now" else "medium",
                "context": f"Roadmap item in {item.get('horizon')} horizon has no assigned owner",
                "generation": {
                    "type": "derived",
                    "source": "missing_owner",
                    "generated_at": datetime.now().isoformat(),
                },
                "derivation": {
                    "pattern": "missing_owner",
                    "evidence": [{
                        "item": item["name"],
                        "horizon": item.get("horizon")
                    }]
                }
            })

    return derived_questions


def generate_llm_questions(context: dict) -> list[dict]:
    """
    Use Claude to generate questions based on full context.
    """
    import anthropic

    # Build context summary for prompt
    roadmap_summary = "\n".join([
        f"- {item['name']} ({item.get('horizon', 'unknown')})"
        for item in context.get("roadmap_items", [])[:15]
    ])

    arch_summary = "No architecture assessments available"
    if isinstance(context.get("arch_assessments"), list) and context["arch_assessments"]:
        arch_summary = "\n".join([
            f"- Assessment {i+1}: Has {len(a.get('gaps', []))} gaps identified"
            for i, a in enumerate(context["arch_assessments"][:3])
        ])

    comp_summary = "No competitive assessments available"
    if isinstance(context.get("competitive_assessments"), list) and context["competitive_assessments"]:
        comp_summary = "\n".join([
            f"- {a.get('competitor_development', {}).get('competitor', 'Unknown')}: {a.get('headline', 'N/A')[:60]}..."
            for a in context["competitive_assessments"][:3]
        ])

    decisions_summary = "\n".join([
        f"- {d.get('decision', 'Unknown')[:60]}..."
        for d in context.get("active_decisions", [])[:10]
    ]) or "No decisions recorded yet"

    answered_summary = "\n".join([
        f"- {q.get('question', 'Unknown')[:60]}..."
        for q in context.get("answered_questions", [])[:10]
    ]) or "No questions answered yet"

    pending_summary = "\n".join([
        f"- [{q.get('priority', 'medium')}] {q.get('question', 'Unknown')[:60]}..."
        for q in context.get("pending_questions", [])[:15]
    ]) or "No pending questions"

    prompt = f"""You are analyzing a product roadmap and related materials to identify questions that need answers.

## CURRENT ROADMAP
{roadmap_summary if roadmap_summary else "No roadmap items available"}

## ARCHITECTURE ASSESSMENT FINDINGS
{arch_summary}

## COMPETITIVE INTELLIGENCE
{comp_summary}

## DECISIONS ALREADY MADE
{decisions_summary}

## QUESTIONS ALREADY ANSWERED
{answered_summary}

## QUESTIONS ALREADY PENDING
{pending_summary}

---

## YOUR TASK

Generate questions that need answers to move this roadmap forward.

IMPORTANT:
- Do NOT generate questions that are already answered above
- Do NOT generate questions that are already pending above
- Do NOT generate questions that are resolved by decisions above
- Focus on NEW questions surfaced by the current state

For each question, specify:
- question: The question text
- audience: engineering | leadership | product
- category: feasibility | investment | direction | trade-off | alignment | timing | scope | dependency | ownership | validation
- priority: critical | high | medium | low
- context: Why this question arose (reference specific roadmap items, assessments, or sources)

Return as JSON array:
```json
[
    {{
        "question": "...",
        "audience": "...",
        "category": "...",
        "priority": "...",
        "context": "..."
    }}
]
```

Generate 3-8 questions depending on complexity. Quality over quantity.
"""

    # Check for API key
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        st.warning("⚠️ ANTHROPIC_API_KEY not set. Skipping LLM question generation. Only derived questions will be generated.")
        return []

    try:
        import httpx
        import json
        import re

        # Create client with SSL verification disabled for development
        http_client = httpx.Client(verify=False)

        try:
            client = anthropic.Anthropic(api_key=api_key, http_client=http_client)

            message = client.messages.create(
                model="claude-sonnet-4-5-20250929",  # Updated to current model
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}]
            )

            response_text = message.content[0].text

            # Extract JSON from response
            json_match = re.search(r'\[[\s\S]*\]', response_text)
            if json_match:
                questions = json.loads(json_match.group())
            else:
                questions = json.loads(response_text)

            # Add generation metadata
            for q in questions:
                q["generation"] = {
                    "type": "llm",
                    "source": "holistic_analysis",
                    "generated_at": datetime.now().isoformat(),
                }

            return questions
        finally:
            # Always close the http client
            http_client.close()

    except Exception as e:
        st.error(f"❌ Error generating LLM questions: {str(e)[:200]}")
        st.info("💡 Continuing with derived questions only. Check API key and network connection.")
        return []


def deduplicate_questions(
    new_questions: list[dict],
    existing_questions: list[dict],
    similarity_threshold: float = 0.85
) -> tuple[list[dict], list[dict]]:
    """
    Deduplicate new questions against existing ones.

    Returns:
        - unique_new: Questions that are genuinely new
        - duplicates: Questions that match existing ones (with match info)
    """

    unique_new = []
    duplicates = []

    # Get embeddings for existing questions
    existing_embeddings = {}
    for eq in existing_questions:
        if eq.get("status") != "obsolete":
            eq_text = eq.get("question", "")
            eq_embedding = get_embedding(eq_text)
            if eq_embedding:
                existing_embeddings[eq.get("id", "unknown")] = {
                    "question": eq,
                    "embedding": eq_embedding
                }

    for new_q in new_questions:
        new_embedding = get_embedding(new_q.get("question", ""))
        if not new_embedding:
            continue

        # Check against all existing
        is_duplicate = False
        best_match = None
        best_similarity = 0

        for eq_id, eq_data in existing_embeddings.items():
            similarity = cosine_similarity(new_embedding, eq_data["embedding"])

            if similarity > best_similarity:
                best_similarity = similarity
                best_match = eq_data["question"]

            if similarity >= similarity_threshold:
                is_duplicate = True
                break

        if is_duplicate:
            duplicates.append({
                "new_question": new_q,
                "matches": best_match,
                "similarity": best_similarity
            })
        else:
            # Add deduplication metadata
            new_q["deduplication"] = {
                "fingerprint": "_".join(extract_key_terms_simple(new_q.get("question", "").lower())[:5]),
                "similar_questions": [best_match.get("id")] if best_match and best_similarity > 0.6 else []
            }

            unique_new.append(new_q)

    return unique_new, duplicates


def mark_obsolete_questions(
    existing_questions: list[dict],
    decisions: list[dict]
) -> list[dict]:
    """
    Mark questions as obsolete if they've been resolved by decisions.

    Returns list of questions that were marked obsolete.
    """

    newly_obsolete = []

    for question in existing_questions:
        if question.get("status") in ["obsolete", "answered"]:
            continue

        # Check if any decision resolves this question
        question_embedding = get_embedding(question.get("question", ""))
        if not question_embedding:
            continue

        for decision in decisions:
            if decision.get("status") != "active":
                continue

            decision_text = f"{decision.get('decision', '')} {decision.get('rationale', '')}"
            decision_embedding = get_embedding(decision_text)
            if not decision_embedding:
                continue

            similarity = cosine_similarity(question_embedding, decision_embedding)

            if similarity > 0.75:
                # Decision likely resolves this question
                question["status"] = "obsolete"
                question["lifecycle"] = question.get("lifecycle", {})
                question["lifecycle"]["obsoleted_at"] = datetime.now().isoformat()
                question["lifecycle"]["obsoleted_by"] = decision.get("id")
                question["lifecycle"]["obsolete_reason"] = f"Resolved by decision: {decision.get('decision', '')[:50]}..."

                newly_obsolete.append(question)
                break

    return newly_obsolete


def generate_questions_holistic() -> dict:
    """
    Main function to generate questions holistically.

    Returns generation results summary.
    """

    run_id = f"gen_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    # Step 1: Gather context
    context = gather_generation_context()

    # Step 2: Derive questions from graph
    derived_questions = derive_questions_from_graph(context)

    # Step 3: Generate LLM questions
    llm_questions = generate_llm_questions(context)

    # Step 4: Combine all new questions
    all_new_questions = derived_questions + llm_questions

    # Step 5: Deduplicate against existing
    unique_questions, duplicates = deduplicate_questions(
        all_new_questions,
        context["existing_questions"]
    )

    # Step 6: Mark obsolete questions
    newly_obsolete = mark_obsolete_questions(
        context["existing_questions"],
        context["active_decisions"]
    )

    # Step 7: Assign IDs and finalize
    existing_questions = context["existing_questions"]

    for q in unique_questions:
        q["id"] = f"q_{datetime.now().strftime('%Y%m%d%H%M%S')}_{len(existing_questions)}"
        q["status"] = "pending"
        q["created_at"] = datetime.now().isoformat()
        q["generation"]["generation_run_id"] = run_id
        q["lifecycle"] = {
            "created_at": datetime.now().isoformat(),
            "refreshed_at": datetime.now().isoformat(),
            "times_refreshed": 1
        }
        existing_questions.append(q)

    # Step 8: Update refresh timestamp on preserved questions
    for eq in existing_questions:
        if eq.get("status") == "pending" and eq.get("id") and eq not in unique_questions:
            eq["lifecycle"] = eq.get("lifecycle", {})
            eq["lifecycle"]["refreshed_at"] = datetime.now().isoformat()
            eq["lifecycle"]["times_refreshed"] = eq["lifecycle"].get("times_refreshed", 0) + 1

    # Step 9: Save
    save_questions(existing_questions)

    # Step 10: Return results summary
    results = {
        "questions_generated": len(unique_questions),
        "llm_generated": len([q for q in unique_questions if q.get("generation", {}).get("type") == "llm"]),
        "derived": len([q for q in unique_questions if q.get("generation", {}).get("type") == "derived"]),
        "duplicates_found": len(duplicates),
        "questions_marked_obsolete": len(newly_obsolete),
        "total_pending_after": len([q for q in existing_questions if q.get("status") == "pending"])
    }

    return results


# ========== ASK YOUR ROADMAP ==========

class QueryIntent(Enum):
    """Types of questions users can ask."""
    WHAT = "what"  # What is X? What does X do?
    WHY = "why"  # Why did we choose X? Why is X important?
    HOW = "how"  # How will X work? How do we implement X?
    WHEN = "when"  # When will X happen? When was X decided?
    WHO = "who"  # Who owns X? Who decided X?
    STATUS = "status"  # What's the status of X?
    COMPARISON = "comparison"  # How does X compare to Y?
    DEPENDENCY = "dependency"  # What depends on X? What does X depend on?
    CONFLICT = "conflict"  # Are there conflicts with X?
    GENERAL = "general"  # General exploration


@dataclass
class ParsedQuery:
    """Structured representation of a user query."""
    original: str
    intent: QueryIntent
    topics: List[str]  # CPQ, Catalog, Experiences, etc.
    keywords: List[str]  # Key terms for search
    modifiers: Dict[str, any]  # priority, recency, severity, scope
    time_context: Optional[str]  # now, next, later, future, past


@dataclass
class RetrievalResult:
    """Results from multi-source retrieval."""
    chunks: List[Dict]  # From LanceDB
    chunk_graph_nodes: List[Dict]  # From chunk context graph
    unified_graph_nodes: Dict[str, List[Dict]]  # By node type


def extract_topics(query: str) -> List[str]:
    """Extract topic mentions from query."""
    topics = []

    # Known topic keywords
    topic_keywords = {
        "CPQ": ["cpq", "configure", "pricing", "quote"],
        "Catalog": ["catalog", "catalogue", "products"],
        "Experiences": ["experience", "experiences", "ux", "ui"],
        "Acquisition": ["acquisition", "acquire", "acquiring"],
        "Intelligent Agents": ["agent", "agents", "ai", "intelligent"],
        "Analytics": ["analytics", "analysis", "reporting", "insights"],
    }

    query_lower = query.lower()

    for topic, keywords in topic_keywords.items():
        if any(kw in query_lower for kw in keywords):
            topics.append(topic)

    return topics


def extract_keywords(query: str) -> List[str]:
    """Extract important keywords from query."""
    import re

    # Remove stop words
    stop_words = {
        "what", "why", "how", "when", "where", "who", "which", "is", "are", "was", "were",
        "the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with",
        "by", "from", "up", "about", "into", "through", "during", "before", "after",
        "above", "below", "between", "under", "again", "further", "then", "once", "here",
        "there", "all", "both", "each", "few", "more", "most", "other", "some", "such",
        "no", "nor", "not", "only", "own", "same", "so", "than", "too", "very", "can",
        "will", "just", "should", "now", "do", "does", "did", "has", "have", "had", "been"
    }

    # Tokenize and filter
    words = re.findall(r'\b\w+\b', query.lower())
    keywords = [w for w in words if w not in stop_words and len(w) > 2]

    return keywords[:10]  # Limit to 10 most important


def extract_modifiers(query: str) -> Dict[str, any]:
    """Extract query modifiers (priority, recency, etc.)."""
    modifiers = {}

    query_lower = query.lower()

    # Priority
    if any(w in query_lower for w in ["critical", "urgent", "important", "priority"]):
        modifiers["priority"] = "high"

    # Recency
    if any(w in query_lower for w in ["recent", "latest", "new", "current"]):
        modifiers["recency"] = "recent"

    # Severity
    if any(w in query_lower for w in ["blocker", "blocking", "severe", "major"]):
        modifiers["severity"] = "high"

    # Scope
    if any(w in query_lower for w in ["all", "everything", "comprehensive", "complete"]):
        modifiers["scope"] = "broad"
    elif any(w in query_lower for w in ["specific", "particular", "just", "only"]):
        modifiers["scope"] = "narrow"

    return modifiers


def extract_time_context(query: str) -> Optional[str]:
    """Extract time context from query."""
    query_lower = query.lower()

    if any(w in query_lower for w in ["now", "current", "today", "immediate"]):
        return "now"
    elif any(w in query_lower for w in ["next", "soon", "upcoming", "short-term"]):
        return "next"
    elif any(w in query_lower for w in ["later", "future", "long-term", "eventually"]):
        return "later"
    elif any(w in query_lower for w in ["past", "previous", "before", "earlier", "decided"]):
        return "past"

    return None


def detect_intent(query: str) -> QueryIntent:
    """Detect the intent of the query."""
    query_lower = query.lower()

    # Intent patterns
    if query_lower.startswith("what") or "what is" in query_lower or "what are" in query_lower:
        return QueryIntent.WHAT
    elif query_lower.startswith("why") or "why did" in query_lower or "why is" in query_lower:
        return QueryIntent.WHY
    elif query_lower.startswith("how") or "how do" in query_lower or "how will" in query_lower:
        return QueryIntent.HOW
    elif query_lower.startswith("when") or "when will" in query_lower or "when was" in query_lower:
        return QueryIntent.WHEN
    elif query_lower.startswith("who") or "who is" in query_lower or "who owns" in query_lower:
        return QueryIntent.WHO
    elif any(w in query_lower for w in ["status", "progress", "state", "where are we"]):
        return QueryIntent.STATUS
    elif any(w in query_lower for w in ["compare", "comparison", "vs", "versus", "difference"]):
        return QueryIntent.COMPARISON
    elif any(w in query_lower for w in ["depend", "dependency", "requires", "needs", "prerequisite"]):
        return QueryIntent.DEPENDENCY
    elif any(w in query_lower for w in ["conflict", "contradiction", "disagree", "inconsistent"]):
        return QueryIntent.CONFLICT

    return QueryIntent.GENERAL


def parse_query(query: str) -> ParsedQuery:
    """Parse user query into structured representation."""
    return ParsedQuery(
        original=query,
        intent=detect_intent(query),
        topics=extract_topics(query),
        keywords=extract_keywords(query),
        modifiers=extract_modifiers(query),
        time_context=extract_time_context(query)
    )


def expand_via_chunk_graph(seed_chunk_ids: List[str], max_hops: int = 1) -> List[Dict]:
    """
    Expand retrieval via chunk context graph.
    Uses BFS to traverse chunk relationships.
    """
    try:
        graph = ContextGraph().load()
        if not graph or not graph.graph:
            return []

        expanded = []
        visited = set(seed_chunk_ids)
        queue = [(chunk_id, 0) for chunk_id in seed_chunk_ids]

        while queue:
            current_id, hops = queue.pop(0)

            if hops >= max_hops:
                continue

            # Get neighbors
            if current_id in graph.graph:
                for neighbor_id in graph.graph.neighbors(current_id):
                    if neighbor_id not in visited:
                        visited.add(neighbor_id)

                        # Get node data
                        node_data = graph.graph.nodes.get(neighbor_id, {})
                        expanded.append({
                            "id": neighbor_id,
                            "content": node_data.get("content", ""),
                            "lens": node_data.get("lens", "unknown"),
                            "source_name": node_data.get("source_name", ""),
                            "source_path": node_data.get("source_file", ""),
                            "hops_from_seed": hops + 1
                        })

                        # Add to queue for next hop
                        if hops + 1 < max_hops:
                            queue.append((neighbor_id, hops + 1))

        return expanded
    except Exception as e:
        st.warning(f"Could not expand via chunk graph: {e}")
        return []


def traverse_unified_graph(seed_chunk_ids: List[str], topics: List[str], max_hops: int = 2) -> Dict[str, List[Dict]]:
    """
    Traverse unified knowledge graph from seed chunks.
    Returns nodes organized by type (decision, question, assessment, roadmap_item, gap).
    """
    try:
        graph = UnifiedContextGraph.load()
        if not graph or not graph.graph:
            return {}

        results = {
            "decision": [],
            "answered_question": [],
            "pending_question": [],
            "assessment": [],
            "roadmap_item": [],
            "gap": []
        }

        visited = set(seed_chunk_ids)
        queue = [(chunk_id, 0) for chunk_id in seed_chunk_ids]

        while queue:
            current_id, hops = queue.pop(0)

            if hops >= max_hops:
                continue

            # Get neighbors in unified graph (both directions)
            if current_id in graph.graph:
                # Check outgoing edges (this node -> neighbors)
                for neighbor_id in graph.graph.neighbors(current_id):
                    if neighbor_id not in visited:
                        visited.add(neighbor_id)

                        # Get node type
                        node_data = graph.graph.nodes.get(neighbor_id, {})
                        node_type = node_data.get("node_type", "")  # FIX: was "type", should be "node_type"

                        # Filter by topics if specified
                        if topics:
                            node_content = str(node_data).lower()
                            if not any(topic.lower() in node_content for topic in topics):
                                continue

                        # Add to results by type
                        if node_type == "decision":
                            results["decision"].append(node_data)
                        elif node_type == "question":
                            # Get actual question data from nested 'data' field
                            q_data = node_data.get("data", node_data)
                            if q_data.get("status") == "answered":
                                results["answered_question"].append(q_data)
                            else:
                                results["pending_question"].append(q_data)
                        elif node_type == "assessment":
                            results["assessment"].append(node_data.get("data", node_data))
                        elif node_type == "roadmap_item":
                            results["roadmap_item"].append(node_data.get("data", node_data))
                        elif node_type == "gap":
                            results["gap"].append(node_data.get("data", node_data))

                        # Continue traversal
                        if hops + 1 < max_hops:
                            queue.append((neighbor_id, hops + 1))

                # FIX: Also check incoming edges (predecessors -> this node)
                for predecessor_id in graph.graph.predecessors(current_id):
                    if predecessor_id not in visited:
                        visited.add(predecessor_id)

                        # Get node type
                        node_data = graph.graph.nodes.get(predecessor_id, {})
                        node_type = node_data.get("node_type", "")

                        # Filter by topics if specified
                        if topics:
                            node_content = str(node_data).lower()
                            if not any(topic.lower() in node_content for topic in topics):
                                continue

                        # Add to results by type
                        if node_type == "decision":
                            results["decision"].append(node_data.get("data", node_data))
                        elif node_type == "question":
                            q_data = node_data.get("data", node_data)
                            if q_data.get("status") == "answered":
                                results["answered_question"].append(q_data)
                            else:
                                results["pending_question"].append(q_data)
                        elif node_type == "assessment":
                            results["assessment"].append(node_data.get("data", node_data))
                        elif node_type == "roadmap_item":
                            results["roadmap_item"].append(node_data.get("data", node_data))
                        elif node_type == "gap":
                            results["gap"].append(node_data.get("data", node_data))

                        # Continue traversal
                        if hops + 1 < max_hops:
                            queue.append((predecessor_id, hops + 1))

        return results
    except Exception as e:
        st.warning(f"Could not traverse unified graph: {e}")
        return {}


def retrieve_full_context(parsed_query: ParsedQuery, top_k: int = 20) -> RetrievalResult:
    """
    Main retrieval function that orchestrates multi-source retrieval.

    Process:
    1. Semantic search in LanceDB using query keywords
    2. Expand via chunk context graph (BFS, max_hops=1)
    3. Traverse unified knowledge graph (max_hops=2)
    4. Filter by topics if specified
    """

    # Step 1: Semantic search in LanceDB
    search_query = " ".join(parsed_query.keywords)
    chunks = retrieve_chunks(search_query, top_k=top_k)

    # Extract chunk IDs
    chunk_ids = [c.get("id", "") for c in chunks if c.get("id")]

    # Step 2: Expand via chunk graph
    chunk_graph_nodes = expand_via_chunk_graph(chunk_ids, max_hops=1)

    # Step 3: Traverse unified graph
    unified_graph_nodes = traverse_unified_graph(chunk_ids, parsed_query.topics, max_hops=2)

    # Step 4: Filter chunks by topics if specified
    if parsed_query.topics:
        filtered_chunks = []
        for chunk in chunks:
            chunk_content = chunk.get("content", "").lower()
            if any(topic.lower() in chunk_content for topic in parsed_query.topics):
                filtered_chunks.append(chunk)
        chunks = filtered_chunks

    return RetrievalResult(
        chunks=chunks,
        chunk_graph_nodes=chunk_graph_nodes,
        unified_graph_nodes=unified_graph_nodes
    )


def assemble_context_for_synthesis(retrieval: RetrievalResult) -> str:
    """
    Assemble retrieved context organized by authority hierarchy for Claude synthesis.

    Authority levels (highest to lowest):
    1. Decisions (L1)
    2. Answered Questions (L2)
    3. Assessments (L3)
    4. Roadmap Items (L4)
    5. Gaps (L5)
    6. Chunks (L6)
    7. Pending Questions (L7)
    """

    context_parts = []

    # L1: Decisions
    decisions = retrieval.unified_graph_nodes.get("decision", [])
    if decisions:
        context_parts.append("# DECISIONS (Authority Level 1 - HIGHEST)\n")
        for decision in decisions[:5]:  # Limit to most relevant
            context_parts.append(f"**Decision:** {decision.get('decision', 'N/A')}")
            context_parts.append(f"**Rationale:** {decision.get('rationale', 'N/A')}")
            context_parts.append(f"**Status:** {decision.get('status', 'N/A')}")
            context_parts.append("")

    # L2: Answered Questions
    answered = retrieval.unified_graph_nodes.get("answered_question", [])
    if answered:
        context_parts.append("# ANSWERED QUESTIONS (Authority Level 2)\n")
        for q in answered[:5]:
            context_parts.append(f"**Q:** {q.get('question', 'N/A')}")
            context_parts.append(f"**A:** {q.get('answer', {}).get('answer', 'N/A')}")
            context_parts.append("")

    # L3: Assessments
    assessments = retrieval.unified_graph_nodes.get("assessment", [])
    if assessments:
        context_parts.append("# ASSESSMENTS (Authority Level 3)\n")
        for assessment in assessments[:3]:
            context_parts.append(f"**Assessment:** {assessment.get('summary', 'N/A')[:200]}")
            context_parts.append("")

    # L4: Roadmap Items
    roadmap_items = retrieval.unified_graph_nodes.get("roadmap_item", [])
    if roadmap_items:
        context_parts.append("# ROADMAP ITEMS (Authority Level 4)\n")
        for item in roadmap_items[:5]:
            context_parts.append(f"**Item:** {item.get('name', 'N/A')} ({item.get('horizon', 'N/A')})")
            context_parts.append(f"**Description:** {item.get('description', 'N/A')[:150]}")
            context_parts.append("")

    # L5: Gaps
    gaps = retrieval.unified_graph_nodes.get("gap", [])
    if gaps:
        context_parts.append("# GAPS (Authority Level 5)\n")
        for gap in gaps[:3]:
            context_parts.append(f"**Gap:** {gap.get('description', 'N/A')[:150]}")
            context_parts.append(f"**Severity:** {gap.get('severity', 'N/A')}")
            context_parts.append("")

    # L6: Chunks (organized by lens)
    if retrieval.chunks:
        context_parts.append("# SOURCE DOCUMENTS (Authority Level 6)\n")

        # Organize by lens
        chunks_by_lens = {}
        for chunk in retrieval.chunks[:15]:
            lens = chunk.get("lens", "unknown")
            if lens not in chunks_by_lens:
                chunks_by_lens[lens] = []
            chunks_by_lens[lens].append(chunk)

        # Output by lens priority
        lens_order = ["your-voice", "team-conversational", "team-structured",
                     "business-framework", "engineering", "external-analyst"]

        for lens in lens_order:
            if lens in chunks_by_lens:
                context_parts.append(f"## {lens.upper()}")
                for chunk in chunks_by_lens[lens][:3]:
                    context_parts.append(f"**Source:** {chunk.get('source_name', 'N/A')}")
                    context_parts.append(f"{chunk.get('content', '')[:300]}")
                    context_parts.append("")

    # L7: Pending Questions (lowest authority)
    pending = retrieval.unified_graph_nodes.get("pending_question", [])
    if pending:
        context_parts.append("# PENDING QUESTIONS (Authority Level 7 - LOWEST)\n")
        for q in pending[:3]:
            context_parts.append(f"- {q.get('question', 'N/A')}")

    return "\n".join(context_parts)


def synthesize_answer(parsed_query: ParsedQuery, context: str) -> Dict:
    """
    Use Claude to synthesize an answer from the assembled context.

    Returns:
    {
        "answer": str,
        "confidence": str (high/medium/low),
        "sources_cited": List[str],
        "related_questions": List[str],
        "follow_ups": List[str]
    }
    """
    import anthropic
    import httpx

    # Check for API key
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {
            "answer": "⚠️ ANTHROPIC_API_KEY not set. Cannot generate answer.",
            "confidence": "none",
            "sources_cited": [],
            "related_questions": [],
            "follow_ups": []
        }

    # Build synthesis prompt
    prompt = f"""You are a product roadmap assistant. Answer the user's question based on the provided context.

## USER QUESTION
{parsed_query.original}

## QUERY ANALYSIS
- Intent: {parsed_query.intent.value}
- Topics: {', '.join(parsed_query.topics) if parsed_query.topics else 'General'}
- Time Context: {parsed_query.time_context or 'Not specified'}

## CONTEXT (ORGANIZED BY AUTHORITY)
The context below is organized from highest authority (decisions) to lowest (pending questions).
When sources conflict, prioritize higher authority sources.

{context}

---

## YOUR TASK

Provide a comprehensive answer that:

1. **Directly answers the question** based on the context
2. **Cites specific sources** by authority level (e.g., "According to Decision X..." or "Based on the your-voice document...")
3. **Highlights conflicts** if information contradicts across authority levels
4. **Acknowledges gaps** if the context doesn't fully answer the question
5. **Provides confidence level** (high/medium/low) based on source quality and completeness

## OUTPUT FORMAT

Answer:
[Your comprehensive answer here, with source citations]

Confidence: [high/medium/low]
Reasoning: [Why this confidence level]

Related Pending Questions:
- [Question 1 from pending questions that relates to this topic]
- [Question 2]

Suggested Follow-ups:
- [Follow-up question 1 the user might want to ask]
- [Follow-up question 2]

Keep the answer focused and actionable. Cite sources explicitly."""

    try:
        # Create client with SSL bypass
        http_client = httpx.Client(verify=False)

        try:
            client = anthropic.Anthropic(api_key=api_key, http_client=http_client)

            message = client.messages.create(
                model="claude-sonnet-4-5-20250929",
                max_tokens=3000,
                messages=[{"role": "user", "content": prompt}]
            )

            response_text = message.content[0].text

            # Parse response
            answer_text = response_text
            confidence = "medium"  # Default
            related_questions = []
            follow_ups = []

            # Extract confidence if present
            if "Confidence:" in response_text:
                lines = response_text.split("\n")
                for i, line in enumerate(lines):
                    if line.startswith("Confidence:"):
                        conf_line = line.replace("Confidence:", "").strip().lower()
                        if "high" in conf_line:
                            confidence = "high"
                        elif "low" in conf_line:
                            confidence = "low"
                        else:
                            confidence = "medium"

                        # Extract answer (everything before Confidence)
                        answer_text = "\n".join(lines[:i]).strip()
                        break

            # Extract related questions
            if "Related Pending Questions:" in response_text:
                section = response_text.split("Related Pending Questions:")[1]
                if "Suggested Follow-ups:" in section:
                    section = section.split("Suggested Follow-ups:")[0]

                related_questions = [
                    line.strip("- ").strip()
                    for line in section.split("\n")
                    if line.strip().startswith("-")
                ]

            # Extract follow-ups
            if "Suggested Follow-ups:" in response_text:
                section = response_text.split("Suggested Follow-ups:")[1]
                follow_ups = [
                    line.strip("- ").strip()
                    for line in section.split("\n")
                    if line.strip().startswith("-")
                ]

            return {
                "answer": answer_text,
                "confidence": confidence,
                "sources_cited": [],  # Could parse from answer text if needed
                "related_questions": related_questions[:3],
                "follow_ups": follow_ups[:3]
            }
        finally:
            http_client.close()

    except Exception as e:
        return {
            "answer": f"❌ Error generating answer: {str(e)[:200]}",
            "confidence": "none",
            "sources_cited": [],
            "related_questions": [],
            "follow_ups": []
        }


def ask_roadmap(query: str) -> Dict:
    """
    Main entry point for Ask Your Roadmap feature.

    Process:
    1. Parse query
    2. Retrieve full context
    3. Assemble context by authority
    4. Synthesize answer via Claude
    5. Return structured result
    """

    # Step 1: Parse query
    parsed_query = parse_query(query)

    # Step 2: Retrieve full context
    retrieval = retrieve_full_context(parsed_query, top_k=20)

    # Step 3: Assemble context
    context = assemble_context_for_synthesis(retrieval)

    # Step 4: Synthesize answer
    result = synthesize_answer(parsed_query, context)

    # Step 5: Add metadata
    result["query_analysis"] = {
        "intent": parsed_query.intent.value,
        "topics": parsed_query.topics,
        "keywords": parsed_query.keywords,
        "time_context": parsed_query.time_context
    }

    result["retrieval_stats"] = {
        "chunks": len(retrieval.chunks),
        "chunk_graph_nodes": len(retrieval.chunk_graph_nodes),
        "decisions": len(retrieval.unified_graph_nodes.get("decision", [])),
        "answered_questions": len(retrieval.unified_graph_nodes.get("answered_question", [])),
        "assessments": len(retrieval.unified_graph_nodes.get("assessment", [])),
        "roadmap_items": len(retrieval.unified_graph_nodes.get("roadmap_item", [])),
        "gaps": len(retrieval.unified_graph_nodes.get("gap", [])),
        "pending_questions": len(retrieval.unified_graph_nodes.get("pending_question", []))
    }

    return result


# ========== SAVE Q&A TO OPEN QUESTIONS ==========

def summarize_answer_for_context(answer: str, max_length: int = 500) -> str:
    """
    Summarize the synthesized answer for use as question context.
    Strips markdown and limits length.
    """
    import re

    # Remove markdown formatting
    text = answer
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # Bold
    text = re.sub(r'\*([^*]+)\*', r'\1', text)      # Italic
    text = re.sub(r'\[Source:[^\]]+\]', '', text)   # Source citations
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)  # Links
    text = re.sub(r'#{1,6}\s*', '', text)           # Headers
    text = re.sub(r'\n{2,}', ' ', text)             # Multiple newlines
    text = re.sub(r'\s{2,}', ' ', text)             # Multiple spaces

    text = text.strip()

    # Truncate if needed
    if len(text) > max_length:
        text = text[:max_length - 3] + "..."

    return f"Based on analysis: {text}"


def get_source_display_name(source: dict) -> str:
    """Get a display name for a source reference."""

    full_data = source.get("full_data", {})
    source_type = source.get("type", "chunk")

    if source_type == "decision":
        return full_data.get("decision", "Decision")[:50]
    elif source_type == "gap":
        return full_data.get("description", "Gap")[:50]
    elif source_type == "assessment":
        return f"{full_data.get('assessment_type', 'Unknown').title()} Assessment"
    elif source_type == "chunk":
        return full_data.get("source_name", "Source document")
    else:
        return source.get("id", "Unknown source")


def find_existing_qa_question(query: str) -> Optional[dict]:
    """Check if this query was already saved as a question."""

    questions = load_questions()

    for q in questions:
        # Check if it came from ask_roadmap
        if q.get("generation", {}).get("source") == "ask_roadmap":
            # Check if query matches
            if q.get("qa_session", {}).get("query", "").lower() == query.lower():
                return q

    return None


def save_qa_to_open_questions(
    query: str,
    answer_result: Dict,
    audience: str,
    category: str,
    priority: str,
    include_answer_as_context: bool = True,
    link_sources: bool = True,
    topic_filter: Optional[str] = None
) -> dict:
    """
    Save a Q&A from Ask Your Roadmap to Open Questions.

    Args:
        query: The original question asked
        answer_result: The complete answer result dict from ask_roadmap()
        audience: engineering | leadership | product
        category: Question category
        priority: critical | high | medium | low
        include_answer_as_context: Whether to include the answer as context
        link_sources: Whether to link the source references
        topic_filter: The topic filter that was applied (if any)

    Returns:
        The created question dict
    """

    question_id = f"q_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    # Build context from answer
    context = ""
    if include_answer_as_context:
        # Summarize the answer for context (strip markdown, limit length)
        context = summarize_answer_for_context(answer_result['answer'], max_length=500)

    # Build source references
    source_references = []
    if link_sources:
        # Get sources from retrieval stats
        retrieval_stats = answer_result.get('retrieval_stats', {})

        # Add summary of sources by type
        if retrieval_stats.get('decisions', 0) > 0:
            source_references.append({
                "type": "decision",
                "count": retrieval_stats['decisions'],
                "relevance": "Referenced in answer synthesis"
            })

        if retrieval_stats.get('assessments', 0) > 0:
            source_references.append({
                "type": "assessment",
                "count": retrieval_stats['assessments'],
                "relevance": "Referenced in answer synthesis"
            })

        if retrieval_stats.get('gaps', 0) > 0:
            source_references.append({
                "type": "gap",
                "count": retrieval_stats['gaps'],
                "relevance": "Referenced in answer synthesis"
            })

        if retrieval_stats.get('chunks', 0) > 0:
            source_references.append({
                "type": "chunk",
                "count": retrieval_stats['chunks'],
                "relevance": "Source documents"
            })

    # Create the question
    question = {
        "id": question_id,
        "question": query,
        "audience": audience,
        "category": category,
        "priority": priority,
        "status": "pending",
        "created_at": datetime.now().isoformat(),

        "generation": {
            "type": "user_query",
            "source": "ask_roadmap",
            "generated_at": datetime.now().isoformat(),
        },

        "context": context,

        "synthesized_answer": {
            "answer": answer_result['answer'],
            "confidence": answer_result['confidence'],
            "generated_at": datetime.now().isoformat(),
            "retrieval_stats": answer_result.get('retrieval_stats', {})
        },

        "source_references": source_references,

        "qa_session": {
            "query": query,
            "topic_filter": topic_filter,
            "session_timestamp": datetime.now().isoformat()
        },

        "validation": None
    }

    # Save to questions list
    questions = load_questions()
    questions.append(question)
    save_questions(questions)

    return question


def render_save_to_questions_ui(query: str, answer_result: Dict, topic_filter: Optional[str] = None, unique_id: int = 0):
    """
    Render the UI for saving a Q&A to Open Questions.
    """

    st.divider()

    with st.container(border=True):
        st.markdown("### 📌 Save to Open Questions")
        st.caption("Add this Q&A to your Open Questions list for follow-up and decision-making.")

        # Check if already saved
        existing = find_existing_qa_question(query)
        if existing:
            st.info(f"✅ This question was already saved as **{existing['id']}**")
            col1, col2 = st.columns(2)
            with col1:
                if st.button("View in Open Questions →", key=f"view_existing_qa_{unique_id}"):
                    st.session_state.current_page = "📝 Open Questions"
                    st.session_state.highlight_question = existing['id']
                    st.rerun()
            return

        # Options
        col1, col2, col3 = st.columns(3)

        with col1:
            audience = st.selectbox(
                "Audience",
                ["leadership", "engineering", "product"],
                index=0,
                key=f"save_qa_audience_{unique_id}"
            )

        with col2:
            category = st.selectbox(
                "Category",
                ["direction", "investment", "feasibility", "trade-off",
                 "alignment", "timing", "scope", "dependency"],
                index=0,
                key=f"save_qa_category_{unique_id}"
            )

        with col3:
            priority = st.selectbox(
                "Priority",
                ["high", "critical", "medium", "low"],
                index=0,
                key=f"save_qa_priority_{unique_id}"
            )

        # Include options
        col1, col2 = st.columns(2)

        with col1:
            include_answer = st.checkbox(
                "Include synthesized answer as context",
                value=True,
                key=f"save_qa_include_answer_{unique_id}",
                help="The answer will be saved as context for the question"
            )

        with col2:
            link_sources = st.checkbox(
                "Link source references",
                value=True,
                key=f"save_qa_link_sources_{unique_id}",
                help="Source citations will be linked to the question"
            )

        # Save button
        if st.button("📌 Save to Open Questions", type="primary", key=f"save_qa_button_{unique_id}"):
            saved_question = save_qa_to_open_questions(
                query=query,
                answer_result=answer_result,
                audience=audience,
                category=category,
                priority=priority,
                include_answer_as_context=include_answer,
                link_sources=link_sources,
                topic_filter=topic_filter
            )

            st.success(f"✅ Saved as **{saved_question['id']}**")

            col1, col2 = st.columns(2)
            with col1:
                if st.button("View in Open Questions →", key=f"view_saved_qa_{unique_id}"):
                    st.session_state.current_page = "📝 Open Questions"
                    st.session_state.highlight_question = saved_question['id']
                    st.rerun()
            with col2:
                if st.button("Make Decision Now →", key=f"decide_saved_qa_{unique_id}"):
                    st.session_state.current_page = "📝 Open Questions"
                    st.session_state.answering_question_id = saved_question['id']
                    st.rerun()


def render_qa_synthesized_answer(synthesized_answer: dict):
    """Render the synthesized answer from a Q&A session."""

    st.markdown("**📝 Synthesized Answer (from Q&A)**")

    confidence = synthesized_answer.get("confidence", "medium")
    confidence_icons = {"high": "🟢", "medium": "🟡", "low": "🟠", "none": "⚪"}
    st.caption(f"Confidence: {confidence_icons.get(confidence, '⚪')} {confidence.title()}")

    # Show answer (collapsed for space)
    answer_text = synthesized_answer.get("answer", "No answer recorded")

    # Truncate for display
    if len(answer_text) > 500:
        st.markdown(answer_text[:500] + "...")
        with st.expander("Show full answer"):
            st.markdown(answer_text)
    else:
        st.markdown(answer_text)

    # Stats
    stats = synthesized_answer.get("retrieval_stats", {})
    if stats:
        st.caption(
            f"Generated from: "
            f"Decisions: {stats.get('decisions', 0)} | "
            f"Assessments: {stats.get('assessments', 0)} | "
            f"Gaps: {stats.get('gaps', 0)} | "
            f"Chunks: {stats.get('chunks', 0)}"
        )


# ========== GRAPH DIAGNOSTICS ==========

def diagnose_graph_contents():
    """Print diagnostic information about the unified graph."""

    graph = UnifiedContextGraph.load()

    if not graph:
        print("❌ ERROR: Unified graph not loaded")
        return None

    print("=" * 60)
    print("UNIFIED GRAPH DIAGNOSTICS")
    print("=" * 60)

    # 1. Check node counts by type
    print("\n1. NODE COUNTS BY TYPE:")
    print("-" * 40)

    node_types = {}
    for node_id in graph.graph.nodes():
        node_data = graph.graph.nodes[node_id]
        node_type = node_data.get("node_type", "unknown")
        node_types[node_type] = node_types.get(node_type, 0) + 1

    for node_type, count in sorted(node_types.items()):
        status = "✅" if count > 0 else "❌"
        print(f"  {status} {node_type}: {count}")

    # 2. Check node_indices
    print("\n2. NODE INDICES:")
    print("-" * 40)

    if hasattr(graph, 'node_indices'):
        for node_type, nodes in graph.node_indices.items():
            print(f"  {node_type}: {len(nodes)} indexed")
    else:
        print("  ❌ graph.node_indices not found!")

    # 3. Check edge types
    print("\n3. EDGE TYPES:")
    print("-" * 40)

    edge_types = {}
    for u, v, data in graph.graph.edges(data=True):
        edge_type = data.get("edge_type", "unknown")
        edge_types[edge_type] = edge_types.get(edge_type, 0) + 1

    for edge_type, count in sorted(edge_types.items()):
        print(f"  {edge_type}: {count}")

    # 4. Check connectivity from chunks to other types
    print("\n4. CHUNK CONNECTIVITY:")
    print("-" * 40)

    chunk_ids = [n for n, d in graph.graph.nodes(data=True) if d.get("node_type") == "chunk"]

    connections_to_types = {
        "decision": 0,
        "roadmap_item": 0,
        "assessment": 0,
        "gap": 0,
        "question": 0
    }

    for chunk_id in chunk_ids[:100]:  # Sample first 100
        for neighbor in graph.graph.neighbors(chunk_id):
            neighbor_type = graph.graph.nodes[neighbor].get("node_type", "unknown")
            if neighbor_type in connections_to_types:
                connections_to_types[neighbor_type] += 1

        for predecessor in graph.graph.predecessors(chunk_id):
            pred_type = graph.graph.nodes[predecessor].get("node_type", "unknown")
            if pred_type in connections_to_types:
                connections_to_types[pred_type] += 1

    for target_type, count in connections_to_types.items():
        status = "✅" if count > 0 else "❌"
        print(f"  {status} Chunks connected to {target_type}: {count}")

    # 5. Sample some non-chunk nodes
    print("\n5. SAMPLE NON-CHUNK NODES:")
    print("-" * 40)

    for node_type in ["decision", "roadmap_item", "assessment", "gap"]:
        nodes = [(n, d) for n, d in graph.graph.nodes(data=True) if d.get("node_type") == node_type]
        if nodes:
            node_id, node_data = nodes[0]
            print(f"\n  Sample {node_type}:")
            print(f"    ID: {node_id}")
            print(f"    Data keys: {list(node_data.keys())}")

            # Check edges
            out_edges = list(graph.graph.neighbors(node_id))
            in_edges = list(graph.graph.predecessors(node_id))
            print(f"    Outgoing edges: {len(out_edges)}")
            print(f"    Incoming edges: {len(in_edges)}")
        else:
            print(f"\n  ❌ No {node_type} nodes found!")

    print("\n" + "=" * 60)

    return node_types, edge_types, connections_to_types


# ========== SOURCES PAGE COMPONENTS ==========

def render_lens_distribution():
    """Render visual distribution of sources by lens."""

    st.subheader("📊 Source Distribution by Lens")

    # Get counts by lens
    lens_counts = get_chunks_by_lens_count()

    if not lens_counts:
        st.info("No sources ingested yet")
        return

    # Define lens colors and order
    lens_order = [
        "your-voice",
        "team-structured",
        "team-conversational",
        "business-framework",
        "engineering",
        "sales",
        "external-analyst"
    ]

    lens_colors = {
        "your-voice": "#e74c3c",
        "team-structured": "#3498db",
        "team-conversational": "#2ecc71",
        "business-framework": "#9b59b6",
        "engineering": "#f39c12",
        "sales": "#1abc9c",
        "external-analyst": "#34495e"
    }

    max_count = max(lens_counts.values()) if lens_counts else 1

    for lens in lens_order:
        count = lens_counts.get(lens, 0)
        if count == 0:
            continue

        pct = count / max_count
        color = lens_colors.get(lens, "#95a5a6")

        col1, col2, col3 = st.columns([2, 4, 1])

        with col1:
            st.markdown(f"**{lens}**")

        with col2:
            # Create a simple progress bar
            st.progress(pct)

        with col3:
            st.caption(f"{count} chunks")


def get_chunks_by_lens_count() -> dict:
    """Get count of chunks grouped by lens."""

    try:
        # Try LanceDB table access
        db = init_db()
        table = db.open_table("roadmap_chunks")
        df = table.to_pandas()
        if not df.empty:
            return df.groupby("lens").size().to_dict()
    except:
        pass

    # Fallback: count from graph
    try:
        graph = UnifiedContextGraph.load()
        if graph:
            lens_counts = {}
            for chunk in graph.node_indices.get("chunk", {}).values():
                lens = chunk.get("lens") if isinstance(chunk, dict) else getattr(chunk, 'lens', "unknown")
                lens_counts[lens] = lens_counts.get(lens, 0) + 1
            return lens_counts
    except:
        pass

    return {}


# ========== PAGE: INGEST MATERIALS ==========

def page_ingest():
    st.title("📥 Ingest Materials")

    # Lens distribution
    render_lens_distribution()

    st.divider()

    st.markdown("Upload documents and tag them with authority levels (lenses)")

    # File upload method
    st.subheader("Upload Files")

    lens = st.selectbox(
        "Select Lens (Authority Level)",
        options=VALID_LENSES,
        help="Choose the authority level for these documents"
    )

    # Lens descriptions
    lens_help = {
        "your-voice": "Your strategic vision (highest authority)",
        "team-structured": "Formal team documentation",
        "team-conversational": "Meeting notes, discussions",
        "business-framework": "OKRs, business strategy",
        "engineering": "Technical docs (veto power on feasibility)",
        "external-analyst": "Market research (lowest authority)"
    }
    st.info(f"**{lens}**: {lens_help.get(lens, '')}")

    # Chunking method selection
    st.subheader("Chunking Method")
    chunking_method = st.radio(
        "Select chunking strategy",
        ["Agentic (Claude-powered)", "Structure-aware (fallback)"],
        index=0,
        help="""
        - **Agentic**: Highest quality, uses Claude to find logical boundaries (~$0.01-0.05 per document)
        - **Structure-aware**: Simple token-based splitting, no LLM cost
        """
    )

    use_agentic = chunking_method == "Agentic (Claude-powered)"

    if use_agentic:
        st.info("⚠️ Agentic chunking calls Claude API. Estimated cost: ~$0.02-0.05 per document.")

    uploaded_files = st.file_uploader(
        "Choose files",
        accept_multiple_files=True,
        type=['pdf', 'docx', 'pptx', 'xlsx', 'csv', 'txt', 'md', 'json']
    )

    if st.button("📥 Ingest Uploaded Files", disabled=not uploaded_files):
        if not uploaded_files:
            st.warning("Please upload at least one file")
            return

        progress_bar = st.progress(0)
        status_text = st.empty()

        success_count = 0
        error_count = 0

        for i, uploaded_file in enumerate(uploaded_files):
            try:
                status_text.text(f"Processing {uploaded_file.name}...")

                # Save file to materials directory
                lens_dir = MATERIALS_DIR / lens
                lens_dir.mkdir(parents=True, exist_ok=True)
                file_path = lens_dir / uploaded_file.name
                file_path.write_bytes(uploaded_file.read())

                # Parse and index
                text = parse_document(file_path)
                if not text.strip():
                    st.warning(f"⚠️ {uploaded_file.name}: Empty document, skipping")
                    error_count += 1
                    continue

                chunks = chunk_with_fallback(text, str(file_path), lens, use_agentic=use_agentic)
                if not chunks:
                    st.warning(f"⚠️ {uploaded_file.name}: No chunks generated, skipping")
                    error_count += 1
                    continue

                index_chunks(chunks, str(file_path))
                st.success(f"✓ {uploaded_file.name} ({len(chunks)} chunks)")
                success_count += 1

            except Exception as e:
                st.error(f"✗ {uploaded_file.name}: {str(e)}")
                error_count += 1

            progress_bar.progress((i + 1) / len(uploaded_files))

        status_text.text(f"Complete! {success_count} succeeded, {error_count} failed")
        st.session_state.index_stats = None  # Clear cache

        # Rebuild context graph
        if success_count > 0:
            with st.spinner("Rebuilding context graph..."):
                graph_stats = rebuild_context_graph()
                if graph_stats:
                    st.success(f"✓ Context graph updated: {graph_stats['nodes']} nodes, {graph_stats['edges']} edges")

    # Alternative: Path-based ingestion
    st.divider()
    st.subheader("Ingest from Path")

    col1, col2 = st.columns([3, 1])
    with col1:
        folder_path = st.text_input("Folder Path", placeholder="/path/to/documents")
    with col2:
        path_lens = st.selectbox("Lens", options=VALID_LENSES, key="path_lens")

    if st.button("📂 Ingest from Path"):
        if not folder_path:
            st.warning("Please provide a path")
            return

        path_obj = Path(folder_path)
        if not path_obj.exists():
            st.error(f"Path not found: {folder_path}")
            return

        files = [f for f in path_obj.rglob("*") if f.is_file() and not f.name.startswith('.')]

        if not files:
            st.warning("No files found in path")
            return

        st.info(f"Found {len(files)} files")

        progress_bar = st.progress(0)
        success_count = 0

        for i, file in enumerate(files):
            try:
                text = parse_document(file)
                if text.strip():
                    chunks = chunk_with_fallback(text, str(file), path_lens, use_agentic=use_agentic)
                    if chunks:
                        index_chunks(chunks, str(file))
                        success_count += 1
            except Exception:
                pass

            progress_bar.progress((i + 1) / len(files))

        st.success(f"Ingested {success_count}/{len(files)} files successfully!")
        st.session_state.index_stats = None

        # Rebuild context graph
        if success_count > 0:
            with st.spinner("Rebuilding context graph..."):
                graph_stats = rebuild_context_graph()
                if graph_stats:
                    st.success(f"✓ Context graph updated: {graph_stats['nodes']} nodes, {graph_stats['edges']} edges")


# ========== PAGE: MANAGE MATERIALS ==========

def page_manage():
    st.title("📁 Manage Materials")
    st.markdown("View, organize, and manage your uploaded documents")

    # Get all materials
    materials = get_all_materials()

    if not materials:
        st.warning("⚠️ No materials found. Upload documents in **Ingest Materials** first.")
        return

    st.success(f"✓ Found {len(materials)} files across {len(set(m['lens'] for m in materials))} lenses")

    # Filter by lens
    st.subheader("Filter & View")
    col1, col2 = st.columns([3, 1])
    with col1:
        selected_lens = st.multiselect(
            "Filter by Lens",
            options=["All"] + VALID_LENSES,
            default=["All"]
        )
    with col2:
        if st.button("🔄 Refresh"):
            st.rerun()

    # Filter materials
    if "All" not in selected_lens and selected_lens:
        filtered_materials = [m for m in materials if m['lens'] in selected_lens]
    else:
        filtered_materials = materials

    # Display as dataframe with controls
    st.subheader(f"Materials ({len(filtered_materials)} files)")

    for material in filtered_materials:
        with st.expander(f"📄 {material['file']} - [{material['lens']}]"):
            col1, col2, col3 = st.columns([2, 2, 1])

            with col1:
                st.write(f"**Lens:** {material['lens']}")
                st.write(f"**Size:** {material['size_mb']}")
                st.write(f"**Modified:** {material['modified'].strftime('%Y-%m-%d %H:%M')}")

            with col2:
                # Move to different lens
                new_lens = st.selectbox(
                    "Move to lens:",
                    options=[l for l in VALID_LENSES if l != material['lens']],
                    key=f"move_{material['path']}"
                )

                if st.button("🔄 Move", key=f"btn_move_{material['path']}"):
                    if move_file_to_lens(material['path'], new_lens):
                        st.success(f"✓ Moved to {new_lens}")
                        st.info("⚠️ Note: You'll need to re-index this file for changes to take effect")
                        st.rerun()

            with col3:
                if st.button("🗑️ Delete", key=f"btn_delete_{material['path']}", type="secondary"):
                    if st.checkbox(f"Confirm delete?", key=f"confirm_{material['path']}"):
                        if delete_material_file(material['path']):
                            st.success("✓ Deleted")
                            st.rerun()

            # File path info
            st.caption(f"Path: {material['path']}")

    # Bulk operations
    st.divider()
    st.subheader("Bulk Operations")

    col1, col2 = st.columns(2)

    with col1:
        st.write("**Re-index All Materials**")
        st.caption("Clears the index and re-ingests all materials. Useful after moving files.")

        if st.button("🔄 Re-index All", type="primary"):
            if st.checkbox("⚠️ This will clear the current index and re-ingest all materials. Continue?"):
                try:
                    # Clear index
                    clear_index()

                    # Re-ingest all materials
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    success_count = 0

                    for i, material in enumerate(materials):
                        try:
                            status_text.text(f"Processing {material['file']}...")
                            text = parse_document(Path(material['path']))
                            if text.strip():
                                chunks = chunk_with_fallback(text, material['path'], material['lens'], use_agentic=True)
                                if chunks:
                                    index_chunks(chunks, material['path'])
                                    success_count += 1
                        except Exception:
                            pass

                        progress_bar.progress((i + 1) / len(materials))

                    status_text.text(f"Complete! Re-indexed {success_count}/{len(materials)} files")
                    st.session_state.index_stats = None
                    st.success("✅ Re-indexing complete!")

                    # Rebuild context graph
                    if success_count > 0:
                        with st.spinner("Rebuilding context graph..."):
                            graph_stats = rebuild_context_graph()
                            if graph_stats:
                                st.success(f"✓ Context graph updated: {graph_stats['nodes']} nodes, {graph_stats['edges']} edges")

                except Exception as e:
                    st.error(f"Error during re-indexing: {e}")

    with col2:
        st.write("**Delete All in Lens**")
        st.caption("Remove all files from a specific lens folder.")

        delete_lens = st.selectbox("Select lens to clear:", options=VALID_LENSES, key="delete_lens")

        if st.button("🗑️ Delete All", type="secondary"):
            if st.checkbox(f"⚠️ This will delete ALL files in {delete_lens}. This cannot be undone!"):
                lens_materials = [m for m in materials if m['lens'] == delete_lens]
                deleted_count = 0

                for material in lens_materials:
                    if delete_material_file(material['path']):
                        deleted_count += 1

                st.success(f"✅ Deleted {deleted_count} files from {delete_lens}")
                st.rerun()


# ========== PAGE: VIEW CHUNKS ==========

def page_chunks():
    st.title("🔍 View Chunks")
    st.markdown("Inspect how documents were chunked and indexed")

    # Explanation section
    with st.expander("ℹ️ How Chunking Works", expanded=False):
        st.markdown("""
        ### Chunking Process

        When you ingest documents, they are broken down into smaller "chunks" for better semantic search and retrieval.

        **Process:**
        1. **Document Parsing**: Files are converted to plain text using the `unstructured` library
        2. **Tokenization**: Text is converted to tokens using tiktoken (cl100k_base encoding, Claude-compatible)
        3. **Chunking**: Text is split into chunks of ~512 tokens each
        4. **Overlap**: Each chunk overlaps with the next by 50 tokens to preserve context across boundaries
        5. **Embedding**: Each chunk is embedded using Voyage AI (voyage-3-large model, 1024 dimensions)
        6. **Indexing**: Embeddings are stored in LanceDB with metadata (lens, source file, token count, etc.)

        **Why Chunk?**
        - Enables precise semantic search over specific sections
        - Fits within context windows for both embedding and LLM models
        - Preserves context through overlap
        - Allows lens-based authority weighting

        **Parameters:**
        - **Chunk Size**: 512 tokens (~350-400 words)
        - **Overlap**: 50 tokens (~35-40 words)
        - **Encoding**: cl100k_base (same as Claude)
        """)

    st.divider()

    # Get chunks
    chunks_df = get_all_chunks()

    if chunks_df is None:
        st.warning("⚠️ No chunks indexed yet. Ingest documents to see chunks here.")
        return

    st.success(f"✓ Found {len(chunks_df)} chunks in the index")

    # Summary statistics
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Total Chunks", f"{len(chunks_df):,}")
    with col2:
        st.metric("Total Tokens", f"{chunks_df['token_count'].sum():,}")
    with col3:
        st.metric("Unique Sources", chunks_df['source_file'].nunique())
    with col4:
        avg_tokens = int(chunks_df['token_count'].mean())
        st.metric("Avg Tokens/Chunk", avg_tokens)

    # Filters
    st.subheader("Filter Chunks")
    col1, col2, col3 = st.columns(3)

    with col1:
        selected_lenses = st.multiselect(
            "Filter by Lens",
            options=["All"] + list(chunks_df['lens'].unique()),
            default=["All"]
        )

    with col2:
        # Get unique source files
        unique_sources = chunks_df['source_file'].unique()
        source_names = [Path(s).name for s in unique_sources]
        selected_source = st.selectbox(
            "Filter by Source File",
            options=["All"] + source_names
        )

    with col3:
        sort_by = st.selectbox(
            "Sort by",
            options=["Created (newest)", "Created (oldest)", "Token Count (high)", "Token Count (low)", "Chunk Index"]
        )

    # Apply filters
    filtered_df = chunks_df.copy()

    if "All" not in selected_lenses and selected_lenses:
        filtered_df = filtered_df[filtered_df['lens'].isin(selected_lenses)]

    if selected_source != "All":
        # Find the full path that matches the selected name
        matching_path = [p for p in unique_sources if Path(p).name == selected_source][0]
        filtered_df = filtered_df[filtered_df['source_file'] == matching_path]

    # Apply sorting
    if sort_by == "Created (newest)":
        filtered_df = filtered_df.sort_values('created_at', ascending=False)
    elif sort_by == "Created (oldest)":
        filtered_df = filtered_df.sort_values('created_at', ascending=True)
    elif sort_by == "Token Count (high)":
        filtered_df = filtered_df.sort_values('token_count', ascending=False)
    elif sort_by == "Token Count (low)":
        filtered_df = filtered_df.sort_values('token_count', ascending=True)
    elif sort_by == "Chunk Index":
        filtered_df = filtered_df.sort_values(['source_file', 'chunk_index'])

    st.info(f"Showing {len(filtered_df)} of {len(chunks_df)} chunks")

    # Display chunks
    st.subheader("Chunk Details")

    # Pagination
    chunks_per_page = 10
    total_pages = (len(filtered_df) + chunks_per_page - 1) // chunks_per_page

    if total_pages > 1:
        page_num = st.number_input("Page", min_value=1, max_value=total_pages, value=1)
        start_idx = (page_num - 1) * chunks_per_page
        end_idx = start_idx + chunks_per_page
        page_df = filtered_df.iloc[start_idx:end_idx]
    else:
        page_df = filtered_df

    # Display each chunk
    for idx, row in page_df.iterrows():
        source_name = Path(row['source_file']).name

        with st.expander(f"🔹 Chunk {row['chunk_index']} from {source_name} [{row['lens']}]"):
            col1, col2 = st.columns([2, 1])

            with col1:
                st.markdown("**Chunk Content:**")
                st.text_area(
                    "Content",
                    value=row['content'],
                    height=200,
                    key=f"chunk_{row['id']}",
                    label_visibility="collapsed"
                )

            with col2:
                st.markdown("**Metadata:**")
                st.write(f"**Lens:** {row['lens']}")
                st.write(f"**Token Count:** {row['token_count']}")
                st.write(f"**Chunk Index:** {row['chunk_index']}")
                st.write(f"**Created:** {row['created_at_str']}")
                st.caption(f"**Source:** {row['source_file']}")
                st.caption(f"**ID:** {row['id']}")

    # Export option
    st.divider()
    st.subheader("Export Chunks")

    col1, col2 = st.columns(2)

    with col1:
        # Export filtered chunks as CSV
        if st.button("📥 Export Filtered as CSV"):
            csv = filtered_df[['id', 'content', 'lens', 'source_file', 'chunk_index', 'token_count', 'created_at_str']].to_csv(index=False)
            st.download_button(
                label="Download CSV",
                data=csv,
                file_name="chunks_export.csv",
                mime="text/csv"
            )

    with col2:
        # Export as JSON
        if st.button("📥 Export Filtered as JSON"):
            json_data = filtered_df[['id', 'content', 'lens', 'source_file', 'chunk_index', 'token_count', 'created_at_str']].to_json(orient='records', indent=2)
            st.download_button(
                label="Download JSON",
                data=json_data,
                file_name="chunks_export.json",
                mime="application/json"
            )


# ========== PAGE: CHUNKING AUDIT ==========

def page_chunking_audit():
    st.title("🔍 Chunking Audit Log")
    st.markdown("Review chunking quality and verification results")

    log_path = DATA_DIR / "chunking_log.jsonl"

    if not log_path.exists():
        st.info("No chunking logs yet. Ingest some documents to see audit information.")
        st.markdown("""
        **What is this page?**

        This page shows detailed information about how documents were chunked:
        - Which chunking method was used (agentic vs structure-aware)
        - Verification results for agentic chunking
        - Any issues detected during chunking
        - Preview of generated chunks
        """)
        return

    # Load logs
    import json
    logs = []
    try:
        with open(log_path) as f:
            for line in f:
                logs.append(json.loads(line))
    except Exception as e:
        st.error(f"Error loading logs: {e}")
        return

    # Summary metrics
    col1, col2, col3, col4 = st.columns(4)
    total_docs = len(logs)
    agentic_docs = len([l for l in logs if l['method'].startswith('agentic')])
    fallback_docs = len([l for l in logs if 'fallback' in l['method'] or l['method'].startswith('structure')])
    issues_docs = len([l for l in logs if not l['verification']['all_valid']])

    col1.metric("Total Documents", total_docs)
    col2.metric("Agentic Chunking", agentic_docs)
    col3.metric("Fallback Used", fallback_docs)
    col4.metric("With Issues", issues_docs)

    # Filters
    st.subheader("Filters")
    col1, col2 = st.columns(2)

    with col1:
        show_only_issues = st.checkbox("Show only documents with issues")

    with col2:
        method_filter = st.selectbox(
            "Filter by method",
            ["All", "Agentic only", "Fallback only"]
        )

    # Apply filters
    filtered_logs = logs

    if show_only_issues:
        filtered_logs = [l for l in filtered_logs if not l['verification']['all_valid']]

    if method_filter == "Agentic only":
        filtered_logs = [l for l in filtered_logs if l['method'].startswith('agentic')]
    elif method_filter == "Fallback only":
        filtered_logs = [l for l in filtered_logs if 'fallback' in l['method'] or l['method'].startswith('structure')]

    # Show results
    st.subheader(f"Results ({len(filtered_logs)} documents)")

    for log in reversed(filtered_logs):  # Most recent first
        status = "✅" if log['verification']['all_valid'] else "⚠️"
        source_name = Path(log['source_path']).name

        with st.expander(f"{status} {source_name} — {log['chunk_count']} chunks — {log['method']}"):
            col1, col2, col3 = st.columns(3)
            col1.write(f"**Lens:** {log['lens']}")
            col2.write(f"**Method:** {log['method']}")
            col3.write(f"**Time:** {log['timestamp'][:19]}")

            if not log['verification']['all_valid']:
                st.error(f"❌ Issues detected in {len(log['verification']['issues'])} chunks")

                for issue in log['verification']['issues']:
                    st.write(f"**Chunk {issue['chunk_index']}:**")
                    for problem in issue['issues']:
                        st.write(f"- {problem.get('severity', 'UNKNOWN')}: {problem.get('message', 'No message')}")
            else:
                st.success(f"✅ All {log['verification']['valid_count']} chunks verified successfully")

            # Show sample chunks
            st.write("**Sample Chunks:**")
            for chunk in log['chunks_preview']:
                st.code(f"[Chunk {chunk['index']}] {chunk['content_length']} chars\n{chunk['content_preview']}...")

                if chunk.get('entities'):
                    st.caption(f"Entities: {', '.join(chunk['entities'])}")


# ========== PAGE: CONTEXT GRAPH ==========

def render_authority_result_card(item: dict, item_type: str):
    """Render a single result card based on type."""

    data = item.get("data", item)
    similarity = item.get("similarity", 0)

    # Check if superseded
    is_superseded = item.get("superseded", False) or data.get("superseded_by")

    # Card styling
    if is_superseded:
        st.markdown(
            f"""<div style="opacity: 0.6; border-left: 3px solid #ff6b6b; padding-left: 10px;">""",
            unsafe_allow_html=True
        )

    with st.container(border=True):
        # Header with similarity score
        col1, col2 = st.columns([4, 1])

        with col1:
            if item_type == "decision":
                st.markdown(f"**{data.get('id', 'Unknown')}**")
                st.write(data.get("decision", ""))

            elif item_type == "chunk":
                source = data.get("source_file", "Unknown")
                if "/" in source:
                    source = source.split("/")[-1]
                lens = data.get("lens", "unknown")
                st.markdown(f"**{source}** ({lens})")
                content = data.get("content", data.get("text", ""))[:200]
                st.text(content + ("..." if len(content) >= 200 else ""))

            elif item_type == "roadmap_item":
                st.markdown(f"**{data.get('name', 'Unknown')}** ({data.get('horizon', '')})")
                st.write(data.get("description", "")[:150])

            elif item_type == "gap":
                severity = data.get("severity", "unknown")
                severity_icon = {"critical": "🔴", "significant": "🟠", "moderate": "🟡", "minor": "🟢"}.get(severity, "⚪")
                st.markdown(f"{severity_icon} **{data.get('description', '')[:100]}**")

            elif item_type == "assessment":
                st.markdown(f"**{data.get('assessment_type', '').title()} Assessment**")
                st.write(data.get("summary", "")[:150])

            elif item_type in ["answered_question", "pending_question"]:
                st.markdown(f"**{data.get('question', '')[:100]}**")
                if data.get("answer"):
                    st.caption(f"Answer: {data['answer'][:100]}...")

        with col2:
            st.caption(f"Score: {similarity:.2f}")

        # Superseded indicator
        if is_superseded:
            superseded_by = data.get("superseded_by") or item.get("superseded_by")
            st.warning(f"⚠️ SUPERSEDED by {superseded_by}")

    if is_superseded:
        st.markdown("</div>", unsafe_allow_html=True)


def render_graph_query_results(results: dict):
    """Render graph query results grouped by authority level."""

    st.subheader("Results by Authority Level")

    st.caption("""
    Results are ordered by authority. Decisions (Level 1) override conflicting
    chunks (Level 6). Content marked as superseded should be treated as outdated.
    """)

    # Level 1: Decisions
    decisions = results.get("decisions", [])
    if decisions:
        with st.expander(f"🔴 **LEVEL 1: DECISIONS** ({len(decisions)}) — Highest Authority", expanded=True):
            st.caption("These decisions are final and override conflicting source content.")
            for item in decisions[:10]:
                render_authority_result_card(item, "decision")

    # Level 2: Answered Questions
    answered = results.get("answered_questions", [])
    if answered:
        with st.expander(f"🟠 **LEVEL 2: ANSWERED QUESTIONS** ({len(answered)})", expanded=False):
            st.caption("Questions that have been resolved by decisions.")
            for item in answered[:10]:
                render_authority_result_card(item, "answered_question")

    # Level 3: Assessments
    assessments = results.get("assessments", [])
    if assessments:
        with st.expander(f"🟡 **LEVEL 3: ASSESSMENTS** ({len(assessments)})", expanded=False):
            st.caption("Analyzed intelligence from architecture, competitive, and impact assessments.")
            for item in assessments[:10]:
                render_authority_result_card(item, "assessment")

    # Level 4: Roadmap Items
    roadmap_items = results.get("roadmap_items", [])
    if roadmap_items:
        with st.expander(f"🟢 **LEVEL 4: ROADMAP ITEMS** ({len(roadmap_items)})", expanded=False):
            st.caption("Current roadmap items.")
            for item in roadmap_items[:10]:
                render_authority_result_card(item, "roadmap_item")

    # Level 5: Gaps
    gaps = results.get("gaps", [])
    if gaps:
        with st.expander(f"🔵 **LEVEL 5: GAPS** ({len(gaps)})", expanded=False):
            st.caption("Identified gaps from assessments, not yet addressed by decisions.")
            for item in gaps[:10]:
                render_authority_result_card(item, "gap")

    # Level 6: Chunks
    chunks = results.get("chunks", [])
    if chunks:
        superseded_count = len([c for c in chunks if c.get("superseded")])
        with st.expander(f"⚪ **LEVEL 6: SOURCE CHUNKS** ({len(chunks)}, {superseded_count} superseded)", expanded=False):
            st.caption("Raw source content. Superseded chunks have been overridden by decisions.")
            for item in chunks[:15]:
                render_authority_result_card(item, "chunk")

    # Level 7: Pending Questions
    pending = results.get("pending_questions", [])
    if pending:
        with st.expander(f"⬜ **LEVEL 7: PENDING QUESTIONS** ({len(pending)}) — Lowest Authority", expanded=False):
            st.caption("Open questions not yet resolved.")
            for item in pending[:10]:
                render_authority_result_card(item, "pending_question")


def page_context_graph():
    st.title("🕸️ Context Graph")
    st.markdown("Explore chunk relationships and semantic connections")

    # Add authority-aware query section at the top
    st.subheader("🔍 Authority-Aware Query")
    st.markdown("Search across all knowledge with automatic authority hierarchy")

    col1, col2 = st.columns([3, 1])
    with col1:
        query = st.text_input("Enter your query", placeholder="e.g., What are the dependencies for Q1 features?")
    with col2:
        include_superseded = st.checkbox("Include superseded content", value=False)

    if query and st.button("🔍 Search", type="primary"):
        with st.spinner("Searching unified knowledge graph..."):
            try:
                graph = UnifiedContextGraph.load()
                if graph and graph.graph.number_of_nodes() > 0:
                    results = retrieve_with_authority(query, graph, top_k=20)
                    render_graph_query_results(results)
                else:
                    st.warning("Unified graph is empty. Run 'Sync Graph' from the Dashboard.")
            except Exception as e:
                st.error(f"Error querying graph: {e}")

    st.divider()

    st.markdown("### Legacy Chunk-Level Exploration")

    # Load graph
    try:
        graph = ContextGraph().load()
        stats = graph.get_stats()
    except Exception as e:
        st.warning("⚠️ No context graph found. Ingest materials to build the graph.")
        st.info("The context graph is automatically built when you ingest documents.")
        return

    # Display statistics
    st.subheader("Graph Statistics")
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Nodes (Chunks)", stats["nodes"])
    col2.metric("Edges (Connections)", stats["edges"])
    col3.metric("Components", stats["components"])
    col4.metric("Density", f"{stats['density']:.4f}")

    # Edge type breakdown
    st.subheader("Edge Types")
    if stats["edge_types"]:
        edge_df = pd.DataFrame([
            {"Edge Type": edge_type, "Count": count}
            for edge_type, count in stats["edge_types"].items()
        ])
        st.bar_chart(edge_df.set_index("Edge Type"))

        # Show legend
        with st.expander("📖 Edge Type Descriptions"):
            st.markdown("""
            - **SIMILAR_TO**: Chunks with high semantic similarity (cosine > 0.80)
            - **SAME_SOURCE**: Chunks from the same source document
            - **SAME_LENS**: Chunks with the same authority lens
            - **TOPIC_OVERLAP**: Chunks sharing 2+ key terms
            - **TEMPORAL_OVERLAP**: Chunks referencing the same time periods
            - **SEQUENTIAL**: Adjacent chunks in the same document
            """)
    else:
        st.info("No edges found in the graph.")

    # Interactive exploration
    st.subheader("Explore Connections")

    if stats["nodes"] == 0:
        st.info("No nodes in the graph.")
        return

    # Select a chunk to explore
    chunk_ids = list(graph.graph.nodes())

    # Get node info for display
    node_options = {}
    for chunk_id in chunk_ids[:100]:  # Limit to first 100 for performance
        node_data = graph.graph.nodes[chunk_id]
        source_name = node_data.get('source_name', 'Unknown')
        lens = node_data.get('lens', 'unknown')
        preview = node_data.get('content_preview', '')[:50]
        display_name = f"{source_name} [{lens}] - {preview}..."
        node_options[display_name] = chunk_id

    if not node_options:
        st.info("No chunks available to explore.")
        return

    selected_display = st.selectbox("Select a chunk to explore", list(node_options.keys()))
    selected = node_options[selected_display]

    if selected:
        node_data = graph.graph.nodes[selected]

        # Show chunk info
        st.write("**Chunk Information:**")
        col1, col2, col3 = st.columns(3)
        col1.metric("Lens", node_data.get('lens', 'N/A'))
        col2.metric("Source", node_data.get('source_name', 'N/A'))
        col3.metric("Token Count", node_data.get('token_count', 'N/A'))

        st.write("**Content Preview:**")
        st.text_area("", node_data.get('content_preview', 'N/A'), height=150, disabled=True)

        # Show key terms and time references
        col1, col2 = st.columns(2)
        with col1:
            st.write("**Key Terms:**")
            terms = node_data.get('key_terms', [])
            if terms:
                st.write(", ".join(terms))
            else:
                st.write("None detected")

        with col2:
            st.write("**Time References:**")
            times = node_data.get('time_references', [])
            if times:
                st.write(", ".join(times))
            else:
                st.write("None detected")

        # Show connections
        neighbors = list(graph.graph.neighbors(selected))
        st.write(f"**Connected to {len(neighbors)} other chunks:**")

        if neighbors:
            # Group by edge type
            by_type = {}
            for neighbor in neighbors:
                edge_data = graph.graph.edges[selected, neighbor]
                edge_type = edge_data.get('type', 'UNKNOWN')
                if edge_type not in by_type:
                    by_type[edge_type] = []
                by_type[edge_type].append((neighbor, edge_data))

            # Display by type
            for edge_type, connections in by_type.items():
                with st.expander(f"{edge_type} ({len(connections)} connections)"):
                    for neighbor, edge_data in connections[:5]:  # Show first 5
                        neighbor_data = graph.graph.nodes[neighbor]
                        st.markdown(f"**→ {neighbor_data.get('source_name', neighbor)}**")
                        st.write(f"- **Lens:** {neighbor_data.get('lens', 'N/A')}")
                        st.write(f"- **Weight:** {edge_data.get('weight', 'N/A')}")

                        # Show shared terms/periods if available
                        if 'shared_terms' in edge_data:
                            st.write(f"- **Shared terms:** {', '.join(edge_data['shared_terms'])}")
                        if 'shared_periods' in edge_data:
                            st.write(f"- **Shared periods:** {', '.join(edge_data['shared_periods'])}")

                        st.write(f"- **Preview:** {neighbor_data.get('content_preview', 'N/A')[:100]}...")
                        st.markdown("---")
        else:
            st.info("This chunk has no connections in the graph.")


# ========== PAGE: GENERATE ROADMAP ==========

def page_generate():
    st.title("🔧 Generate Roadmap")
    st.markdown("Generate master roadmap from indexed materials")

    # Check if we have materials
    stats = get_index_stats()
    if not stats:
        st.warning("⚠️ No materials indexed. Please ingest documents first.")
        return

    st.success(f"✓ Ready to generate ({stats['total_chunks']} chunks indexed)")

    # Show existing roadmap if available
    master_path = OUTPUT_DIR / "master_roadmap.md"
    if master_path.exists():
        st.info(f"📄 Existing master roadmap found (last modified: {datetime.fromtimestamp(master_path.stat().st_mtime).strftime('%Y-%m-%d %H:%M')})")

        with st.expander("Preview Existing Roadmap"):
            st.markdown(master_path.read_text())

    # Generation options
    st.subheader("Generation Options")

    additional_context = st.text_area(
        "Additional Context (optional)",
        placeholder="E.g., Focus on Q2 priorities, emphasize platform investments...",
        help="Provide additional instructions for the roadmap generation"
    )

    # Generate button
    if st.button("🚀 Generate Master Roadmap", type="primary"):
        # Check API keys first
        if not os.getenv("ANTHROPIC_API_KEY"):
            st.error("❌ ANTHROPIC_API_KEY not set. Please configure in Settings.")
            return
        if not os.getenv("VOYAGE_API_KEY"):
            st.error("❌ VOYAGE_API_KEY not set. Please configure in Settings.")
            return

        try:
            with st.spinner("Generating roadmap... This may take 30-60 seconds"):
                query = "Generate a comprehensive product roadmap"
                if additional_context:
                    query += f". Additional context: {additional_context}"

                roadmap = generate_roadmap(query)

                st.success("✅ Master roadmap generated successfully!")

                # Display roadmap
                st.subheader("Generated Roadmap")
                st.markdown(roadmap)

                # Download button
                st.download_button(
                    label="📥 Download Roadmap",
                    data=roadmap,
                    file_name="master_roadmap.md",
                    mime="text/markdown"
                )

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            st.error(f"Error generating roadmap: {e}")
            with st.expander("🔍 Error Details"):
                st.code(error_details)

            # Provide helpful suggestions
            st.info("""
            **Troubleshooting:**
            - Check that your API keys are valid in Settings
            - Verify you have internet connectivity
            - Try testing the connection in Settings page
            - Check if the Anthropic API is accessible from your network
            """)


# ========== PAGE: FORMAT BY PERSONA ==========

def page_format():
    st.title("👥 Format by Persona")
    st.markdown("Transform master roadmap for specific audiences")

    # Check if master roadmap exists
    master_path = OUTPUT_DIR / "master_roadmap.md"
    if not master_path.exists():
        st.warning("⚠️ Master roadmap not found. Please generate it first.")
        return

    st.success(f"✓ Master roadmap found (last modified: {datetime.fromtimestamp(master_path.stat().st_mtime).strftime('%Y-%m-%d %H:%M')})")

    # Persona selection
    st.subheader("Select Persona")

    persona = st.radio(
        "Target Audience",
        options=["executive", "product", "engineering"],
        format_func=lambda x: {
            "executive": "📊 Executive (C-level, SVPs) - Strategic view, 500-800 words",
            "product": "📋 Product (PMs, Product Leaders) - Detailed priorities, 2-3 pages",
            "engineering": "⚙️ Engineering (Engineers, Architects) - Technical details, 5-10 pages"
        }[x],
        help="Choose the audience for the formatted roadmap"
    )

    # Check if formatted version exists
    persona_path = OUTPUT_DIR / f"{persona}_roadmap.md"
    if persona_path.exists():
        st.info(f"📄 Existing {persona} roadmap found (last modified: {datetime.fromtimestamp(persona_path.stat().st_mtime).strftime('%Y-%m-%d %H:%M')})")

    # Generate button
    if st.button(f"🎯 Generate {persona.capitalize()} Roadmap", type="primary"):
        try:
            with st.spinner(f"Formatting roadmap for {persona} audience..."):
                formatted = format_for_persona(persona)

                st.success(f"✅ {persona.capitalize()} roadmap generated successfully!")

                # Display with tabs
                tab1, tab2 = st.tabs(["📄 Formatted Roadmap", "📋 Master Roadmap"])

                with tab1:
                    st.markdown(formatted)
                    st.download_button(
                        label=f"📥 Download {persona.capitalize()} Roadmap",
                        data=formatted,
                        file_name=f"{persona}_roadmap.md",
                        mime="text/markdown",
                        key="download_formatted"
                    )

                with tab2:
                    master_text = master_path.read_text()
                    st.markdown(master_text)

        except Exception as e:
            st.error(f"Error formatting roadmap: {e}")

    # Show existing formatted roadmaps
    st.divider()
    st.subheader("Existing Roadmaps")

    for p in ["executive", "product", "engineering"]:
        p_path = OUTPUT_DIR / f"{p}_roadmap.md"
        if p_path.exists():
            with st.expander(f"📄 {p.capitalize()} Roadmap"):
                st.markdown(p_path.read_text())
                st.download_button(
                    label=f"📥 Download",
                    data=p_path.read_text(),
                    file_name=f"{p}_roadmap.md",
                    mime="text/markdown",
                    key=f"download_{p}"
                )


# ========== PAGE: ASK QUESTIONS ==========

def page_ask():
    st.title("💬 Ask Your Roadmap")
    st.markdown("Conversational Q&A powered by multi-source retrieval and Claude synthesis")

    # Check if we have materials
    stats = get_index_stats()
    if not stats:
        st.warning("⚠️ No materials indexed. Please ingest documents first.")
        return

    # Initialize session state for ask history
    if 'ask_history' not in st.session_state:
        st.session_state.ask_history = []

    # Clear history button
    col1, col2 = st.columns([6, 1])
    with col2:
        if st.button("🗑️ Clear", key="clear_ask_history"):
            st.session_state.ask_history = []
            st.rerun()

    # Display chat history
    for idx, msg in enumerate(st.session_state.ask_history):
        if msg['role'] == 'user':
            with st.chat_message("user"):
                st.write(msg['content'])
        else:
            with st.chat_message("assistant"):
                # Display answer
                st.markdown(msg['answer'])

                # Display metadata in columns
                col1, col2, col3 = st.columns(3)

                with col1:
                    confidence = msg.get('confidence', 'medium')
                    conf_emoji = {"high": "🟢", "medium": "🟡", "low": "🔴"}.get(confidence, "⚪")
                    st.caption(f"Confidence: {conf_emoji} {confidence.title()}")

                with col2:
                    intent = msg.get('query_analysis', {}).get('intent', 'general')
                    st.caption(f"Intent: {intent}")

                with col3:
                    topics = msg.get('query_analysis', {}).get('topics', [])
                    if topics:
                        st.caption(f"Topics: {', '.join(topics)}")

                # Retrieval stats
                stats = msg.get('retrieval_stats', {})
                if stats:
                    with st.expander("📊 Retrieval Statistics"):
                        col1, col2, col3, col4 = st.columns(4)
                        col1.metric("Chunks", stats.get('chunks', 0))
                        col2.metric("Decisions", stats.get('decisions', 0))
                        col3.metric("Assessments", stats.get('assessments', 0))
                        col4.metric("Roadmap Items", stats.get('roadmap_items', 0))

                # Related questions
                related = msg.get('related_questions', [])
                if related:
                    with st.expander("🔗 Related Pending Questions"):
                        for q in related:
                            st.markdown(f"- {q}")

                # Follow-ups
                follow_ups = msg.get('follow_ups', [])
                if follow_ups:
                    with st.expander("💡 Suggested Follow-ups"):
                        for fu in follow_ups:
                            if st.button(fu, key=f"followup_{idx}_{hash(fu)}"):
                                # Add as new question
                                st.session_state.next_question = fu
                                st.rerun()

            # === ADD: Save to Open Questions UI ===
            # Show save UI for each Q&A (outside chat_message for better layout)
            if msg.get('role') == 'assistant' and msg.get('query'):
                # Create a unique key for this save UI using index
                render_save_to_questions_ui(
                    query=msg['query'],
                    answer_result=msg,
                    topic_filter=msg.get('query_analysis', {}).get('topics', [None])[0] if msg.get('query_analysis', {}).get('topics') else None,
                    unique_id=idx
                )

    # Question input
    question = st.chat_input("Ask a question about your roadmap...")

    # Handle follow-up question from button
    if 'next_question' in st.session_state:
        question = st.session_state.next_question
        del st.session_state.next_question

    if question:
        # Add user message
        st.session_state.ask_history.append({
            'role': 'user',
            'content': question
        })

        # Get answer using ask_roadmap
        try:
            with st.spinner("🔍 Analyzing query and retrieving context..."):
                result = ask_roadmap(question)

            # Add assistant message with full result for saving
            st.session_state.ask_history.append({
                'role': 'assistant',
                'query': question,  # Store the query for save UI
                'answer': result['answer'],
                'confidence': result['confidence'],
                'query_analysis': result['query_analysis'],
                'retrieval_stats': result['retrieval_stats'],
                'related_questions': result.get('related_questions', []),
                'follow_ups': result.get('follow_ups', [])
            })

            st.rerun()

        except Exception as e:
            st.error(f"❌ Error generating answer: {str(e)[:300]}")
            # Show more details in expander
            with st.expander("Error Details"):
                st.code(str(e))


# ========== PAGE: OPEN QUESTIONS ==========

def page_open_questions():
    st.title("📝 Open Questions")
    st.markdown("Track open questions, submit answers, and manage the decision log")

    # Validation stats
    render_validation_stats()

    st.divider()

    # === GENERATE QUESTIONS SECTION ===
    st.subheader("🔄 Generate Questions")

    # Show what will be analyzed
    try:
        context = gather_generation_context()

        with st.expander("What will be analyzed", expanded=False):
            col1, col2 = st.columns(2)

            with col1:
                st.markdown("**Roadmap & Assessments**")
                st.caption(f"• {len(context.get('roadmap_items', []))} roadmap items")
                st.caption(f"• {len(context.get('arch_assessments', []))} architecture assessments")
                st.caption(f"• {len(context.get('competitive_assessments', []))} competitive assessments")

            with col2:
                st.markdown("**Decisions & Questions**")
                st.caption(f"• {len(context.get('active_decisions', []))} active decisions")
                st.caption(f"• {len(context.get('answered_questions', []))} answered questions")
                st.caption(f"• {len(context.get('pending_questions', []))} pending questions")

        # Generate button
        col1, col2 = st.columns([2, 1])

        with col1:
            if st.button("🔄 Generate Questions", type="primary", use_container_width=True):
                with st.spinner("Analyzing context and generating questions..."):
                    results = generate_questions_holistic()

                st.success(f"✅ Generation complete!")

                # Show results
                rcol1, rcol2, rcol3, rcol4 = st.columns(4)
                rcol1.metric("New Questions", results["questions_generated"])
                rcol2.metric("LLM Generated", results["llm_generated"])
                rcol3.metric("Derived", results["derived"])
                rcol4.metric("Marked Obsolete", results["questions_marked_obsolete"])

                st.caption(f"Total pending: {results['total_pending_after']}")

                # Invalidate source cache
                invalidate_source_cache()

                st.rerun()

        with col2:
            st.caption("Generates new questions based on current roadmap, assessments, and decisions")

    except Exception as e:
        st.error(f"Error preparing generation: {e}")

    st.divider()

    # Maintenance section
    with st.expander("🔧 Maintenance"):
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Clear Source Cache"):
                invalidate_source_cache()
                st.success("Source cache cleared! Sources will be re-searched on next view.")
                st.rerun()
        with col2:
            st.caption("Clear cached source references. Use after ingesting new documents or syncing the graph.")

    st.divider()

    # Load data
    questions = load_questions()
    answers = load_answers()
    decisions = load_decisions()

    # Create tabs
    tab1, tab2, tab3 = st.tabs(["📋 Pending Questions", "✅ Answer Question", "📜 Decision Log"])

    # ========== TAB 1: PENDING QUESTIONS ==========
    with tab1:
        st.subheader("Pending Questions")

        # Filters
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            audience_filter = st.selectbox("Audience", ["All", "engineering", "leadership", "product"])
        with col2:
            priority_filter = st.selectbox("Priority", ["All", "critical", "high", "medium", "low"])
        with col3:
            category_filter = st.selectbox("Category", ["All", "feasibility", "investment", "direction", "trade-off", "alignment", "timing", "scope", "dependency", "ownership"])
        with col4:
            type_filter = st.selectbox(
                "Source",
                ["All", "user_query", "llm", "derived", "legacy"],
                format_func=lambda x: {
                    "All": "All Sources",
                    "user_query": "💬 From Q&A",
                    "llm": "🤖 LLM Generated",
                    "derived": "🔍 Derived",
                    "legacy": "📝 Legacy"
                }.get(x, x)
            )

        # Filter questions
        pending = [q for q in questions if q.get("status", "pending") == "pending"]

        if audience_filter != "All":
            pending = [q for q in pending if q.get("audience", "") == audience_filter]
        if priority_filter != "All":
            pending = [q for q in pending if q.get("priority", "") == priority_filter]
        if category_filter != "All":
            pending = [q for q in pending if q.get("category", "") == category_filter]
        if type_filter != "All":
            if type_filter == "legacy":
                pending = [q for q in pending if not q.get("generation")]
            else:
                pending = [q for q in pending if q.get("generation", {}).get("type") == type_filter]

        # Summary metrics
        col1, col2, col3, col4 = st.columns(4)
        total_pending = len([q for q in questions if q.get("status", "pending") == "pending"])
        critical_count = len([q for q in pending if q.get("priority", "") == "critical"])
        answered_count = len([q for q in questions if q.get("status", "") == "answered"])
        decisions_count = len(decisions)

        col1.metric("Total Pending", total_pending)
        col2.metric("Critical", critical_count)
        col3.metric("Answered", answered_count)
        col4.metric("Decisions Made", decisions_count)

        st.markdown("---")

        if not pending:
            st.info("No pending questions match your filters.")
        else:
            # Display questions by audience
            for audience in ["engineering", "leadership", "product"]:
                audience_qs = [q for q in pending if q.get("audience", "") == audience]
                if not audience_qs:
                    continue

                st.markdown(f"### {audience.title()} ({len(audience_qs)})")

                # Sort by priority
                priority_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
                sorted_qs = sorted(audience_qs, key=lambda x: priority_order.get(x.get("priority", "low"), 4))

                for q in sorted_qs:
                    priority_emoji = {"critical": "🔴", "high": "🟠", "medium": "🟡", "low": "⚪"}.get(q.get("priority", "medium"), "⚪")

                    # Validation status
                    validation = q.get("validation")
                    if validation and validation.get("validated"):
                        val_icon = "👍" if validation.get("is_accurate") else "👎"
                    else:
                        val_icon = "❓"

                    # Generation type badge
                    generation = q.get("generation", {})
                    gen_type = generation.get("type", "legacy")
                    gen_source = generation.get("source", "")

                    if gen_type == "user_query" and gen_source == "ask_roadmap":
                        type_badge = "💬"
                        has_synthesized_answer = True
                    elif gen_type == "llm":
                        type_badge = "🤖"
                        has_synthesized_answer = False
                    elif gen_type == "derived":
                        type_badge = "🔍"
                        has_synthesized_answer = False
                    else:
                        type_badge = "📝"
                        has_synthesized_answer = False

                    question_preview = q['question'][:70] + "..." if len(q['question']) > 70 else q['question']

                    with st.expander(f"{priority_emoji} {val_icon} {type_badge} {question_preview}"):
                        # Header
                        st.write(f"**Question:** {q['question']}")

                        # Generation type info
                        if gen_type == "user_query" and gen_source == "ask_roadmap":
                            st.success("💬 From Q&A - Asked in Ask Your Roadmap")
                        elif gen_type == "llm":
                            st.info("🤖 Generated by LLM analysis")
                        elif gen_type == "derived":
                            st.warning(f"🔍 Derived from {generation.get('source', 'unknown')} pattern")
                        else:
                            st.caption("📝 Legacy question")

                        st.write(f"**Category:** {q.get('category', 'N/A')}")
                        st.write(f"**Priority:** {q.get('priority', 'medium')}")
                        st.write(f"**Context:** {q.get('context', 'None provided')}")

                        if q.get("related_roadmap_items"):
                            st.write(f"**Affects:** {', '.join(q['related_roadmap_items'])}")

                        # === NEW: Show synthesized answer for Q&A questions ===
                        if has_synthesized_answer and q.get("synthesized_answer"):
                            st.divider()
                            render_qa_synthesized_answer(q["synthesized_answer"])

                        # Show derivation evidence for derived questions
                        if gen_type == "derived":
                            derivation = q.get("derivation", {})
                            evidence = derivation.get("evidence", [])

                            if evidence:
                                with st.expander(f"📊 Derivation Evidence ({len(evidence)} items)"):
                                    for ev in evidence:
                                        if "source_name" in ev:
                                            st.markdown(f"**{ev.get('source_name')}** ({ev.get('lens', 'unknown')})")
                                            st.caption(ev.get("content", "")[:200])
                                        else:
                                            st.json(ev)

                        st.write(f"**Created:** {q.get('created_at', 'Unknown')[:10]}")

                        st.divider()

                        # Source references
                        render_question_source_references(q)

                        st.divider()

                        # Validation
                        render_question_validation(q)

                        st.divider()

                        # Quick actions
                        if has_synthesized_answer:
                            # For Q&A questions, show Re-Ask and Mark Obsolete buttons
                            col1, col2 = st.columns(2)
                            if col1.button("🔄 Re-Ask", key=f"reask_{q['id']}"):
                                st.session_state.current_page = "💬 Ask Your Roadmap"
                                # Store the question to pre-fill
                                if 'ask_history' not in st.session_state:
                                    st.session_state.ask_history = []
                                st.rerun()
                            if col2.button("Mark Obsolete", key=f"obs_{q['id']}"):
                                q["status"] = "obsolete"
                                save_questions(questions)
                                st.success("Question marked obsolete")
                                st.rerun()
                        else:
                            # For other questions, show Defer and Mark Obsolete buttons
                            col1, col2 = st.columns(2)
                            if col1.button("Defer", key=f"def_{q['id']}"):
                                q["status"] = "deferred"
                                save_questions(questions)
                                st.success("Question deferred")
                                st.rerun()
                            if col2.button("Mark Obsolete", key=f"obs_{q['id']}"):
                                q["status"] = "obsolete"
                                save_questions(questions)
                                st.success("Question marked obsolete")
                                st.rerun()

    # ========== TAB 2: ANSWER QUESTION ==========
    with tab2:
        st.subheader("Answer Question")

        # Check if a question was selected from Tab 1
        pre_selected_id = st.session_state.get("answering_question_id")

        if pre_selected_id:
            st.info("📌 Question pre-selected from Tab 1. The dropdown below has been set to your selected question.")

        # Select question to answer
        pending = [q for q in questions if q.get("status", "pending") == "pending"]

        if not pending:
            st.info("No pending questions. Questions will be generated after roadmap synthesis.")
        else:
            question_options = {f"{q['id']}: {q['question'][:60]}...": q['id'] for q in pending}

            # Find the index of the pre-selected question
            default_index = 0
            if pre_selected_id:
                for i, (label, q_id) in enumerate(question_options.items()):
                    if q_id == pre_selected_id:
                        default_index = i
                        break

            selected = st.selectbox(
                "Select Question",
                list(question_options.keys()),
                index=default_index
            )

            # Clear the pre-selection after using it
            if pre_selected_id and 'answering_question_id' in st.session_state:
                del st.session_state.answering_question_id

            if selected:
                q_id = question_options[selected]
                question = next(q for q in pending if q["id"] == q_id)

                st.markdown(f"**Question:** {question['question']}")
                st.markdown(f"**Context:** {question.get('context', 'None')}")
                st.markdown(f"**Audience:** {question.get('audience', 'N/A')} | **Category:** {question.get('category', 'N/A')}")

                if question.get("related_roadmap_items"):
                    st.markdown(f"**Affects:** {', '.join(question['related_roadmap_items'])}")

                st.markdown("---")

                # Answer form (Enhanced with decision flow)
                answer_text = st.text_area(
                    "Your Answer",
                    height=120,
                    key=f"answer_{q_id}",
                    placeholder="Provide your answer to this question..."
                )

                col1, col2 = st.columns(2)
                with col1:
                    answered_by = st.text_input("Answered By", key=f"by_{q_id}")
                with col2:
                    confidence = st.selectbox(
                        "Confidence",
                        ["high", "medium", "low"],
                        index=1,
                        key=f"conf_{q_id}"
                    )

                st.divider()

                # Decision creation option
                create_decision = st.checkbox(
                    "✅ Create decision from this answer",
                    value=True,
                    key=f"dec_{q_id}",
                    help="Automatically create a decision record that will be tracked in the Decision Log and override conflicting source content"
                )

                if create_decision:
                    st.markdown("**Decision Details**")

                    # Pre-populate decision from answer
                    default_decision = answer_text[:150] if answer_text else ""

                    decision_statement = st.text_input(
                        "Decision Statement",
                        value=default_decision,
                        key=f"dtext_{q_id}",
                        help="Concise statement of what was decided (will appear in Decision Log)"
                    )

                    decision_rationale = st.text_area(
                        "Rationale",
                        value=answer_text,
                        height=80,
                        key=f"rat_{q_id}",
                        help="Why this decision was made"
                    )

                    decision_implications = st.text_input(
                        "Implications (comma-separated)",
                        key=f"imp_{q_id}",
                        placeholder="e.g., Timeline shifts to Q3, Need to inform customers"
                    )

                    decision_owner = st.text_input(
                        "Decision Owner",
                        value=answered_by,
                        key=f"own_{q_id}",
                        help="Who is responsible for executing this decision"
                    )

                st.divider()

                # Submit button
                can_submit = bool(answer_text and answered_by)
                if create_decision:
                    can_submit = can_submit and bool(decision_statement)

                if st.button("Submit Answer", type="primary", disabled=not can_submit, key=f"submit_{q_id}"):
                    # Save answer
                    answer_record = {
                        "id": f"ans_{datetime.now().strftime('%Y%m%d%H%M%S')}",
                        "question_id": q_id,
                        "answer": answer_text,
                        "answered_by": answered_by,
                        "answered_at": datetime.now().isoformat(),
                        "confidence": confidence,
                        "follow_up_needed": False,
                        "notes": ""
                    }
                    answers.append(answer_record)
                    save_answers(answers)

                    # Update question status
                    question["status"] = "answered"
                    question["answer"] = answer_text
                    question["answered_by"] = answered_by
                    question["answered_at"] = datetime.now().isoformat()
                    save_questions(questions)

                    # Save decision if requested
                    if create_decision:
                        decision_record = {
                            "id": f"dec_{datetime.now().strftime('%Y%m%d%H%M%S')}",
                            "question_id": q_id,
                            "answer_id": answer_record["id"],
                            "decision": decision_statement,
                            "rationale": decision_rationale,
                            "implications": [i.strip() for i in decision_implications.split(",") if i.strip()],
                            "owner": decision_owner,
                            "review_date": None,
                            "status": "active",
                            "created_at": datetime.now().isoformat()
                        }
                        decisions.append(decision_record)
                        save_decisions(decisions)

                        st.success(f"✅ Answer submitted and decision **{decision_record['id']}** created!")

                        # Prompt to sync graph
                        st.info("💡 Run 'Sync Graph' to integrate this decision into the knowledge graph")
                    else:
                        st.success("✅ Answer submitted!")

                    st.rerun()

    # ========== TAB 3: DECISION LOG ==========
    with tab3:
        st.subheader("Decision Log")

        # Filters
        col1, col2 = st.columns(2)
        with col1:
            show_superseded = st.checkbox("Show superseded decisions")
        with col2:
            date_filter = st.date_input("Since", value=None)

        # Filter decisions
        filtered_decisions = decisions
        if not show_superseded:
            filtered_decisions = [d for d in filtered_decisions if d.get("status", "active") == "active"]
        if date_filter:
            date_str = str(date_filter)
            filtered_decisions = [d for d in filtered_decisions if d.get("created_at", "")[:10] >= date_str]

        # Sort by date
        filtered_decisions = sorted(filtered_decisions, key=lambda x: x.get("created_at", ""), reverse=True)

        st.markdown(f"**Total Decisions:** {len(filtered_decisions)}")
        st.markdown("---")

        if not filtered_decisions:
            st.info("No decisions recorded yet. Answer questions to create decisions.")
        else:
            # Display decisions with overrides
            for dec in filtered_decisions:
                status_icon = {"active": "✅", "superseded": "🔄", "revisiting": "🔍"}.get(dec.get("status", "active"), "?")
                created = dec.get("created_at", "Unknown")[:10]
                decision_preview = dec['decision'][:60] + "..." if len(dec['decision']) > 60 else dec['decision']

                with st.container(border=True):
                    # Header
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        st.markdown(f"### {status_icon} {dec['id']}")
                    with col2:
                        st.caption(created)

                    # Decision statement
                    st.markdown(f"**Decision:** {dec['decision']}")

                    # Rationale
                    if dec.get("rationale"):
                        with st.expander("Rationale"):
                            st.write(dec["rationale"])

                    # Implications
                    if dec.get("implications"):
                        st.markdown("**Implications:**")
                        for imp in dec["implications"]:
                            st.markdown(f"- {imp}")

                    # Link to original question
                    question_id = dec.get("question_id")
                    if question_id:
                        question = next((q for q in questions if q["id"] == question_id), None)
                        if question:
                            st.caption(f"Resolves: {question['question'][:80]}...")

                    # Owner
                    if dec.get("owner"):
                        st.caption(f"Owner: {dec['owner']}")

                    st.divider()

                    # OVERRIDES SECTION - Show what this decision overrides
                    overrides = get_decision_overrides(dec["id"])

                    if overrides:
                        with st.expander(f"⚡ Overrides {len(overrides)} source chunk(s)", expanded=False):
                            st.caption("This decision supersedes the following source content:")

                            for chunk in overrides:
                                source_name = chunk.get("source_file", "Unknown").split("/")[-1] if chunk.get("source_file") else "Unknown"
                                st.markdown(f"**{source_name}** ({chunk.get('lens', 'unknown')})")
                                content = chunk.get("content", chunk.get("text", ""))
                                st.text(content[:200] + ("..." if len(content) > 200 else ""))
                                st.divider()
                    else:
                        st.caption("No conflicting source chunks identified")

                    # Actions
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        if st.button("Mark for Review", key=f"review_{dec['id']}"):
                            dec["status"] = "revisiting"
                            save_decision_update(dec)
                            st.rerun()
                    with col2:
                        if st.button("Supersede", key=f"supersede_{dec['id']}"):
                            dec["status"] = "superseded"
                            save_decision_update(dec)
                            st.rerun()
                    with col3:
                        if dec.get("status") != "active":
                            if st.button("Reactivate", key=f"reactivate_{dec['id']}"):
                                dec["status"] = "active"
                                save_decision_update(dec)
                                st.rerun()

            # Export button
            st.markdown("---")
            if st.button("📥 Export Decision Log"):
                md_lines = [
                    "# Decision Log",
                    f"\nGenerated: {datetime.now().isoformat()}",
                    f"Total Decisions: {len(filtered_decisions)}\n",
                    "---\n"
                ]

                for dec in filtered_decisions:
                    status_emoji = {"active": "✅", "superseded": "🔄", "revisiting": "🔍"}.get(dec.get("status", "active"), "?")
                    created = dec.get("created_at", "Unknown")[:10]

                    md_lines.append(f"## {status_emoji} {dec['id']} ({created})\n")
                    md_lines.append(f"**Decision:** {dec['decision']}\n")

                    if dec.get("rationale"):
                        md_lines.append(f"**Rationale:** {dec['rationale']}\n")

                    if dec.get("implications"):
                        md_lines.append("**Implications:**")
                        for imp in dec["implications"]:
                            md_lines.append(f"- {imp}")
                        md_lines.append("")

                    md_lines.append(f"**Owner:** {dec.get('owner', 'Unassigned')}")
                    md_lines.append(f"**Status:** {dec.get('status', 'active')}")

                    question_id = dec.get("question_id")
                    if question_id:
                        md_lines.append(f"**Question ID:** {question_id}")

                    md_lines.append("\n---\n")

                md_content = "\n".join(md_lines)

                st.download_button(
                    "Download Markdown",
                    md_content,
                    file_name=f"decision_log_{datetime.now().strftime('%Y%m%d')}.md",
                    mime="text/markdown"
                )


# ========== PAGE: ARCHITECTURE ALIGNMENT ==========

def page_architecture_alignment():
    st.title("🏗️ Architecture Alignment Analysis")

    st.markdown("""
    This analysis maps functional roadmap items to technical architecture to identify:
    - Whether architecture supports planned features
    - Required technical changes
    - Technical risks and dependencies
    - Engineering questions that need answers
    """)

    # Create tabs
    tab1, tab2, tab3 = st.tabs(["📊 Analysis", "📄 Architecture Docs", "❓ Engineering Questions"])

    # ========== TAB 1: ANALYSIS ==========
    with tab1:
        st.subheader("Roadmap-Architecture Alignment")

        # Check for existing analysis
        alignment_file = OUTPUT_DIR / "architecture-alignment.json"

        col1, col2 = st.columns([3, 1])
        with col2:
            if st.button("🔄 Run Analysis", type="primary"):
                # Check if roadmap exists
                roadmap_file = OUTPUT_DIR / "master_roadmap.md"
                if not roadmap_file.exists():
                    st.error("Roadmap not found. Generate a roadmap first.")
                else:
                    # Check if documents are selected
                    selected_paths = st.session_state.get('selected_doc_paths', [])
                    if not selected_paths:
                        st.error("⚠️ No documents selected. Go to 'Architecture Docs' tab and select documents.")
                    else:
                        with st.spinner("Loading selected architecture documents..."):
                            arch_docs, metadata = load_architecture_documents(selected_files=selected_paths)

                        if not arch_docs:
                            st.error("No documents could be loaded. Check the Architecture Docs tab for details.")
                        else:
                            with st.spinner(f"Analyzing alignment ({len(arch_docs)} docs, {metadata['total_tokens']:,} tokens)... This may take 1-2 minutes"):
                                roadmap_content = roadmap_file.read_text()
                                roadmap = parse_roadmap_for_analysis(roadmap_content)
                                analysis = generate_architecture_alignment(roadmap, arch_docs, use_opus=True)
                                save_alignment_analysis(analysis)

                                # Add questions to system
                                questions = extract_engineering_questions_from_alignment(analysis)
                                added = add_architecture_questions_to_system(questions)

                            st.success(f"Analysis complete! {added} engineering questions added to Open Questions.")
                            st.rerun()

        if alignment_file.exists():
            analysis = load_alignment_analysis()

            # Summary metrics
            assessments = analysis.get("assessments", [])

            col1, col2, col3, col4 = st.columns(4)
            col1.metric("Items Analyzed", len(assessments))
            col2.metric("Full Support", len([a for a in assessments if a.get("architecture_supports") == "full"]))
            col3.metric("Partial Support", len([a for a in assessments if a.get("architecture_supports") == "partial"]))
            col4.metric("No Support", len([a for a in assessments if a.get("architecture_supports") == "no"]))

            st.markdown("---")

            # Filter
            support_filter = st.multiselect(
                "Filter by support level",
                ["full", "partial", "no", "unknown"],
                default=["partial", "no"]
            )

            # Display assessments
            filtered_assessments = [a for a in assessments if a.get("architecture_supports") in support_filter]

            for assessment in filtered_assessments:
                support_icon = {
                    "full": "✅",
                    "partial": "⚠️",
                    "no": "❌",
                    "unknown": "❓"
                }.get(assessment.get("architecture_supports", "unknown"), "❓")

                with st.expander(f"{support_icon} {assessment.get('roadmap_item', 'Unknown')} ({assessment.get('horizon', 'N/A')})"):
                    st.write(f"**Description:** {assessment.get('roadmap_item_description', 'N/A')}")
                    st.write(f"**Architecture Support:** {assessment.get('architecture_supports', 'N/A')} (confidence: {assessment.get('confidence', 'N/A')})")
                    st.write(f"**Summary:** {assessment.get('summary', 'N/A')}")

                    # Supporting components
                    if assessment.get("supporting_components"):
                        st.write("**Supporting Components:**")
                        for comp in assessment["supporting_components"]:
                            st.write(f"  • {comp}")

                    # Required changes
                    if assessment.get("required_changes"):
                        st.write("**Required Changes:**")
                        for change in assessment["required_changes"]:
                            blocking = "🚫 BLOCKING" if change.get("blocking") else ""
                            st.write(f"  • **{change.get('component', 'N/A')}** ({change.get('change_type', 'N/A')}, {change.get('effort', 'N/A')}) — {change.get('description', 'N/A')} {blocking}")

                    # Technical risks
                    if assessment.get("technical_risks"):
                        st.write("**Technical Risks:**")
                        for risk in assessment["technical_risks"]:
                            severity_color = {"critical": "🔴", "high": "🟠", "medium": "🟡", "low": "🟢"}.get(risk.get("severity", "low"), "⚪")
                            st.write(f"  {severity_color} {risk.get('risk', 'N/A')}")
                            st.write(f"     Mitigation: {risk.get('mitigation', 'N/A')}")

                    # Dependencies
                    deps = assessment.get("dependencies", {})
                    if deps.get("prerequisite_work"):
                        st.write(f"**Prerequisites:** {', '.join(deps['prerequisite_work'])}")
                    if deps.get("enables"):
                        st.write(f"**Enables:** {', '.join(deps['enables'])}")

                    # Questions
                    if assessment.get("questions"):
                        st.write(f"**Engineering Questions:** {len(assessment['questions'])}")
                        for q in assessment["questions"]:
                            st.write(f"  • [{q.get('priority', 'N/A')}] {q.get('question', 'N/A')}")

            # Cross-cutting concerns
            cross_cutting = analysis.get("cross_cutting_concerns", {})

            if cross_cutting:
                st.markdown("---")
                st.subheader("Cross-Cutting Concerns")

                if cross_cutting.get("architectural_gaps"):
                    st.write("**Architectural Gaps:**")
                    for gap in cross_cutting["architectural_gaps"]:
                        st.write(f"  • {gap}")

                if cross_cutting.get("systemic_risks"):
                    st.write("**Systemic Risks:**")
                    for risk in cross_cutting["systemic_risks"]:
                        st.write(f"  • {risk}")

                if cross_cutting.get("recommended_adrs"):
                    st.write("**Recommended ADRs:**")
                    for adr in cross_cutting["recommended_adrs"]:
                        st.write(f"  • {adr}")

                if cross_cutting.get("sequencing_recommendations"):
                    st.write("**Sequencing Recommendations:**")
                    st.write(cross_cutting["sequencing_recommendations"])

            # Export
            st.markdown("---")
            col1, col2 = st.columns(2)
            with col1:
                md_content = format_alignment_report(analysis)
                st.download_button(
                    "📥 Export as Markdown",
                    md_content,
                    file_name="architecture-alignment.md",
                    mime="text/markdown"
                )
            with col2:
                import json
                st.download_button(
                    "📥 Export as JSON",
                    json.dumps(analysis, indent=2),
                    file_name="architecture-alignment.json",
                    mime="application/json"
                )

        else:
            st.info("No analysis found. Click 'Run Analysis' to analyze roadmap against architecture.")

    # ========== TAB 2: ARCHITECTURE DOCS ==========
    with tab2:
        st.subheader("Architecture Documents")

        # Upload section
        st.markdown("### Upload Architecture Documents")

        col1, col2 = st.columns(2)
        with col1:
            doc_type = st.selectbox(
                "Document Type",
                ["architecture", "tech-specs"],
                help="Choose where to save the document"
            )

        uploaded_files = st.file_uploader(
            "Upload architecture documents (.md, .txt, .rst)",
            type=["md", "txt", "rst"],
            accept_multiple_files=True,
            help="Upload architecture documents, technical specs, or design docs"
        )

        if uploaded_files:
            if st.button("💾 Save Documents", type="primary"):
                # Create directory if it doesn't exist
                target_dir = Path(f"materials/engineering/{doc_type}")
                target_dir.mkdir(parents=True, exist_ok=True)

                saved_files = []
                for uploaded_file in uploaded_files:
                    # Save file
                    file_path = target_dir / uploaded_file.name
                    file_path.write_bytes(uploaded_file.getvalue())
                    saved_files.append(uploaded_file.name)

                st.success(f"✓ Successfully saved {len(saved_files)} document(s) to {target_dir}")

                # Show uploaded files
                with st.expander("📄 Uploaded Files", expanded=True):
                    for filename in saved_files:
                        st.write(f"  • {filename}")

                st.info("Documents will be loaded when you run the next analysis")

                # Wait a moment before rerun
                import time
                time.sleep(1)
                st.rerun()

        st.markdown("---")
        st.markdown("### Manage Architecture Files")

        # List files from both directories
        st.markdown("**Architecture Documents:**")
        arch_dir = Path("materials/engineering/architecture")
        if arch_dir.exists():
            arch_files = list(arch_dir.rglob("*"))
            arch_files = [f for f in arch_files if f.is_file() and f.suffix in ['.md', '.txt', '.rst']]

            if arch_files:
                for file_path in sorted(arch_files):
                    col1, col2, col3 = st.columns([3, 1, 1])
                    with col1:
                        st.text(f"📄 {file_path.name}")
                    with col2:
                        file_size = file_path.stat().st_size
                        st.text(f"{file_size / 1024:.1f} KB")
                    with col3:
                        if st.button("🗑️ Delete", key=f"del_arch_{file_path.name}"):
                            file_path.unlink()
                            st.success(f"Deleted {file_path.name}")
                            st.rerun()
            else:
                st.info("No files in architecture/ directory")
        else:
            st.info("architecture/ directory not created yet")

        st.markdown("**Technical Specs:**")
        specs_dir = Path("materials/engineering/tech-specs")
        if specs_dir.exists():
            spec_files = list(specs_dir.rglob("*"))
            spec_files = [f for f in spec_files if f.is_file() and f.suffix in ['.md', '.txt', '.rst']]

            if spec_files:
                for file_path in sorted(spec_files):
                    col1, col2, col3 = st.columns([3, 1, 1])
                    with col1:
                        st.text(f"📄 {file_path.name}")
                    with col2:
                        file_size = file_path.stat().st_size
                        st.text(f"{file_size / 1024:.1f} KB")
                    with col3:
                        if st.button("🗑️ Delete", key=f"del_spec_{file_path.name}"):
                            file_path.unlink()
                            st.success(f"Deleted {file_path.name}")
                            st.rerun()
            else:
                st.info("No files in tech-specs/ directory")
        else:
            st.info("tech-specs/ directory not created yet")

        st.markdown("---")
        st.markdown("### Select Documents for Analysis")
        st.caption("Choose which architecture documents to include (max 200K tokens total)")

        try:
            # Scan available documents
            available_docs = scan_architecture_documents()

            if not available_docs:
                st.warning("⚠️ No architecture documents found in materials/engineering/")
                st.info("Upload documents in the 'Architecture Docs' tab above")
            else:
                st.success(f"📚 Found {len(available_docs)} architecture documents")

                # Initialize session state for document selection
                if 'selected_doc_paths' not in st.session_state:
                    # Auto-select documents that fit in budget (up to 200K tokens)
                    st.session_state.selected_doc_paths = []
                    running_total = 0
                    for doc in available_docs:
                        if not doc['too_large'] and running_total + doc['token_count'] <= 200000:
                            st.session_state.selected_doc_paths.append(doc['path'])
                            running_total += doc['token_count']

                # Document selection UI
                st.markdown("**Select documents to include in analysis:**")

                selected_paths = []
                running_tokens = 0

                for doc in available_docs:
                    col1, col2, col3 = st.columns([3, 2, 1])

                    with col1:
                        is_selected = st.checkbox(
                            doc['title'],
                            value=doc['path'] in st.session_state.selected_doc_paths,
                            key=f"doc_select_{doc['path']}",
                            disabled=doc['too_large']
                        )

                        if is_selected:
                            selected_paths.append(doc['path'])
                            running_tokens += doc['token_count']

                    with col2:
                        token_color = "red" if doc['too_large'] else "green" if doc['token_count'] < 50000 else "orange"
                        st.markdown(f":{token_color}[{doc['token_count']:,} tokens]")

                    with col3:
                        st.caption(f"{doc['file_size'] / 1024:.1f} KB")

                    if doc['too_large']:
                        st.warning(f"⚠️ This document exceeds the 200K token limit per file and cannot be loaded")

                # Update session state
                st.session_state.selected_doc_paths = selected_paths

                # Show running total
                st.markdown("---")
                budget_color = "red" if running_tokens > 200000 else "green"
                st.markdown(f"**Selected:** {len(selected_paths)} documents | **Total Tokens:** :{budget_color}[{running_tokens:,} / 200,000]")

                if running_tokens > 200000:
                    st.error("⚠️ Selected documents exceed the 200K token budget. Deselect some documents.")
                elif running_tokens == 0:
                    st.info("ℹ️ Select at least one document to run analysis")
                else:
                    # Load selected documents
                    with st.spinner("Loading selected documents..."):
                        docs, metadata = load_architecture_documents(selected_files=selected_paths)

                    if metadata['skipped_count'] > 0:
                        with st.expander(f"⚠️ {metadata['skipped_count']} documents were skipped"):
                            for skipped in metadata['skipped_files']:
                                st.warning(f"**{skipped['name']}**: {skipped['reason']}")

                    # Show loaded documents
                    st.success(f"✓ **{len(docs)} documents loaded** ({metadata['total_tokens']:,} tokens)")

                    for doc in docs:
                        with st.expander(f"📄 {doc['title']} ({doc['token_count']:,} tokens)"):
                            st.write(f"**Path:** {doc['path']}")
                            st.write(f"**Type:** {doc['doc_type']}")
                            st.write(f"**Last Updated:** {doc.get('last_updated', 'Unknown')[:10]}")

                            if doc.get('key_components'):
                                st.write(f"**Key Components:** {', '.join(doc['key_components'][:10])}")

        except Exception as e:
            st.error(f"❌ Error scanning documents: {e}")
            import traceback
            st.code(traceback.format_exc())

    # ========== TAB 3: ENGINEERING QUESTIONS ==========
    with tab3:
        st.subheader("Engineering Questions from Architecture Analysis")

        questions = load_questions()
        arch_questions = [q for q in questions if q.get("source") == "architecture_alignment"]

        if not arch_questions:
            st.info("No architecture-related questions yet. Run an analysis to generate questions.")
        else:
            pending = [q for q in arch_questions if q.get("status") == "pending"]
            answered = [q for q in arch_questions if q.get("status") == "answered"]

            st.write(f"**{len(pending)} pending** | {len(answered)} answered")

            for q in sorted(pending, key=lambda x: {"critical": 0, "high": 1, "medium": 2, "low": 3}.get(x.get("priority", "low"), 4)):
                priority_icon = {"critical": "🔴", "high": "🟠", "medium": "🟡", "low": "🟢"}.get(q.get("priority", "medium"), "⚪")

                with st.expander(f"{priority_icon} {q.get('question', 'N/A')[:80]}..."):
                    st.write(f"**Question:** {q.get('question', 'N/A')}")
                    st.write(f"**Category:** {q.get('category', 'N/A')}")
                    st.write(f"**Context:** {q.get('context', 'N/A')}")

                    if q.get("related_component"):
                        st.write(f"**Component:** {q['related_component']}")

                    if q.get("related_roadmap_items"):
                        st.write(f"**Roadmap Items:** {', '.join(q['related_roadmap_items'])}")

                    st.markdown("---")
                    st.info("Go to the 📝 Open Questions page to answer this question")


# ========== PAGE: COMPETITIVE INTELLIGENCE ==========

def page_competitive_intelligence():
    st.title("🎯 Competitive Intelligence")

    st.markdown("""
    ### Analyst Assessment Tool

    This tool generates **objective analyst assessments** of competitor developments against your roadmap.

    **Important:** This produces analyst research notes, not strategy recommendations:
    - The roadmap is the source of truth
    - No new features or approaches are suggested
    - Questions are raised for leadership, not answered
    """)

    # Create tabs
    tab1, tab2, tab3 = st.tabs(["📝 Manage Developments", "📊 Run Assessment", "📄 View Assessments"])

    # ========== TAB 1: MANAGE DEVELOPMENTS ==========
    with tab1:
        st.subheader("Competitor Developments")

        st.markdown("### Add New Development")

        with st.form("add_development"):
            col1, col2 = st.columns(2)

            with col1:
                competitor = st.text_input("Competitor Name*", placeholder="e.g., Competitor X")
                dev_type = st.selectbox(
                    "Development Type*",
                    ["product_launch", "feature", "acquisition", "partnership", "funding", "strategy_shift"]
                )
                title = st.text_input("Title*", placeholder="e.g., AI-Powered Pricing Launch")

            with col2:
                source_url = st.text_input("Source URL*", placeholder="https://...")
                announced_date = st.date_input("Announced Date*")
                description = st.text_area("Description*", placeholder="Full description of the development...")

            submitted = st.form_submit_button("➕ Add Development", type="primary")

            if submitted:
                if not all([competitor, title, description, source_url]):
                    st.error("All fields are required")
                else:
                    try:
                        development = add_competitor_development(
                            competitor=competitor,
                            development_type=dev_type,
                            title=title,
                            description=description,
                            source_url=source_url,
                            announced_date=str(announced_date)
                        )
                        st.success(f"✓ Added development: {development['id']}")
                        st.info("Go to the 'Run Assessment' tab to analyze this development")
                        import time
                        time.sleep(1)
                        st.rerun()
                    except Exception as e:
                        st.error(f"Error adding development: {e}")

        st.markdown("---")
        st.markdown("### Existing Developments")

        developments = load_competitor_developments()

        if not developments:
            st.info("No competitor developments tracked yet. Add one above to get started.")
        else:
            for dev in developments:
                with st.expander(f"🔹 {dev['competitor']} — {dev['title']}", expanded=False):
                    col1, col2, col3 = st.columns(3)

                    with col1:
                        st.write(f"**ID:** {dev['id']}")
                        st.write(f"**Type:** {dev['development_type']}")

                    with col2:
                        st.write(f"**Announced:** {dev['announced_date']}")
                        st.write(f"**Added:** {dev['created_at'][:10]}")

                    with col3:
                        # Check if assessed
                        assessments = load_analyst_assessments()
                        assessed = any(a['development_id'] == dev['id'] for a in assessments)
                        if assessed:
                            st.success("✓ Assessed")
                        else:
                            st.warning("⧗ Not yet assessed")

                    st.write(f"**Description:** {dev['description']}")
                    st.write(f"**Source:** {dev['source_url']}")

    # ========== TAB 2: RUN ASSESSMENT ==========
    with tab2:
        st.subheader("Generate Analyst Assessment")

        st.info("**Note:** Assessments are objective analyst research notes, not strategy recommendations.")

        developments = load_competitor_developments()

        if not developments:
            st.warning("No competitor developments found. Add a development in the 'Manage Developments' tab first.")
        else:
            # Select development
            dev_options = {f"{d['id']} — {d['competitor']}: {d['title']}": d['id'] for d in developments}
            selected_label = st.selectbox("Select Development to Assess", list(dev_options.keys()))
            selected_id = dev_options[selected_label]

            development = get_competitor_development(selected_id)

            if development:
                st.markdown("---")
                st.markdown("### Development Details")
                st.write(f"**Competitor:** {development['competitor']}")
                st.write(f"**Title:** {development['title']}")
                st.write(f"**Type:** {development['development_type']}")
                st.write(f"**Announced:** {development['announced_date']}")
                st.write(f"**Description:** {development['description'][:200]}...")

                st.markdown("---")

                # Check prerequisites
                roadmap_file = OUTPUT_DIR / "master_roadmap.md"
                if not roadmap_file.exists():
                    st.error("⚠️ Roadmap not found. Generate a roadmap first.")
                else:
                    col1, col2 = st.columns([3, 1])

                    with col2:
                        use_opus = st.checkbox("Use Opus (higher quality)", value=True)

                    with col1:
                        if st.button("🔄 Run Analyst Assessment", type="primary"):
                            with st.spinner("Generating analyst assessment... This may take 1-2 minutes"):
                                try:
                                    assessment = generate_analyst_assessment(development, use_opus=use_opus)
                                    analysis = assessment['analysis']

                                    st.success("✓ Assessment Complete!")

                                    # Show summary
                                    st.markdown(f"### {analysis['headline']}")
                                    st.write(analysis['executive_summary'])

                                    st.markdown("---")

                                    col1, col2, col3 = st.columns(3)
                                    col1.metric("Impact", analysis['overall_impact'])
                                    col2.metric("Timeline", analysis['impact_timeline'])
                                    col3.metric("Confidence", analysis['confidence'])

                                    st.write(f"**Roadmap Strengths:** {len(analysis.get('roadmap_strengths', []))}")
                                    st.write(f"**Roadmap Gaps:** {len(analysis.get('roadmap_gaps', []))}")
                                    st.write(f"**Strategic Questions:** {len(analysis.get('strategic_questions', []))}")

                                    st.info(f"Strategic questions have been added to Open Questions for leadership review.")

                                    if assessment.get('validation_warnings'):
                                        with st.expander(f"⚠️ {len(assessment['validation_warnings'])} Validation Warnings"):
                                            for warning in assessment['validation_warnings']:
                                                st.warning(warning)

                                    st.rerun()

                                except Exception as e:
                                    st.error(f"Assessment failed: {e}")
                                    import traceback
                                    st.code(traceback.format_exc())

    # ========== TAB 3: VIEW ASSESSMENTS ==========
    with tab3:
        st.subheader("Analyst Assessments")

        assessments = load_analyst_assessments()

        if not assessments:
            st.info("No assessments generated yet. Run an assessment in the 'Run Assessment' tab.")
        else:
            st.success(f"📊 {len(assessments)} assessment(s) generated")

            for assessment in reversed(assessments):  # Most recent first
                analysis = assessment['analysis']
                development = assessment['development']

                with st.expander(f"📄 {analysis['headline']}", expanded=False):
                    st.write(f"**Assessed:** {assessment['assessed_at'][:10]}")
                    st.write(f"**Competitor:** {development['competitor']}")
                    st.write(f"**Development:** {development['title']}")

                    st.markdown("---")

                    st.markdown("### Executive Summary")
                    st.write(analysis['executive_summary'])

                    st.markdown("### Impact Assessment")
                    col1, col2, col3 = st.columns(3)
                    col1.metric("Impact", analysis['overall_impact'])
                    col2.metric("Timeline", analysis['impact_timeline'])
                    col3.metric("Confidence", analysis['confidence'])

                    # Roadmap Strengths
                    st.markdown("### Roadmap Strengths")
                    for strength in analysis.get('roadmap_strengths', []):
                        st.write(f"**{strength['roadmap_item']}** ({strength['horizon']} horizon)")
                        st.write(f"- Coverage: {strength['coverage_level']}")
                        st.write(f"- Timing: {strength['timing_adequacy']}")
                        st.write(f"- {strength['how_it_addresses']}")

                    # Roadmap Gaps
                    st.markdown("### Roadmap Gaps")
                    for gap in analysis.get('roadmap_gaps', []):
                        severity_color = {"critical": "🔴", "significant": "🟠", "moderate": "🟡", "minor": "🟢"}.get(gap['severity'], "⚪")
                        st.write(f"{severity_color} **{gap['gap_description']}**")
                        st.write(f"- Severity: {gap['severity']}")
                        st.write(f"- Competitor has: {gap['competitor_capability']}")

                    # Strategic Questions
                    st.markdown("### Strategic Questions Raised")
                    for i, q in enumerate(analysis.get('strategic_questions', []), 1):
                        st.write(f"{i}. **[{q['question_type'].upper()}]** {q['question']}")
                        st.write(f"   *{q['context']}*")

                    # Export
                    st.markdown("---")
                    markdown_content = format_analyst_assessment_markdown(assessment)
                    st.download_button(
                        "📥 Export as Markdown",
                        markdown_content,
                        file_name=f"competitive_assessment_{assessment['id']}.md",
                        mime="text/markdown",
                        key=f"export_{assessment['id']}"
                    )


# ========== PAGE: SETTINGS ==========

def page_settings():
    st.title("⚙️ Settings")
    st.markdown("Configure API keys and system settings")

    # API Keys section
    st.subheader("API Keys")

    current_anthropic = os.getenv("ANTHROPIC_API_KEY", "")
    current_voyage = os.getenv("VOYAGE_API_KEY", "")

    # Show masked keys
    if current_anthropic:
        st.info(f"Anthropic API Key: {current_anthropic[:10]}...{current_anthropic[-4:]}")
    if current_voyage:
        st.info(f"Voyage AI Key: {current_voyage[:10]}...{current_voyage[-4:]}")

    # Input new keys
    with st.form("api_keys"):
        anthropic_key = st.text_input(
            "Anthropic API Key",
            value=current_anthropic,
            type="password",
            help="Get from: https://console.anthropic.com/"
        )

        voyage_key = st.text_input(
            "Voyage AI API Key",
            value=current_voyage,
            type="password",
            help="Get from: https://dash.voyageai.com/"
        )

        col1, col2 = st.columns(2)
        with col1:
            save_button = st.form_submit_button("💾 Save Settings", type="primary")
        with col2:
            test_button = st.form_submit_button("🧪 Test Connection")

        if save_button:
            if not anthropic_key or not voyage_key:
                st.error("Both API keys are required")
            else:
                save_env_vars(anthropic_key, voyage_key)
                st.success("✅ API keys saved successfully!")

        if test_button:
            try:
                validate_api_keys()
                st.success("✅ API keys are valid!")
            except Exception as e:
                st.error(f"❌ API key validation failed: {e}")

    # Configuration display
    st.divider()
    st.subheader("Current Configuration")

    config_data = {
        "Setting": ["Materials Directory", "Output Directory", "Data Directory", "Valid Lenses"],
        "Value": [
            str(MATERIALS_DIR),
            str(OUTPUT_DIR),
            str(DATA_DIR),
            ", ".join(VALID_LENSES)
        ]
    }
    st.table(pd.DataFrame(config_data))


# ========== MAIN APP ==========

def main():
    # Sidebar navigation
    st.sidebar.title("🗺️ Roadmap Synth")
    st.sidebar.markdown("---")

    page = st.sidebar.radio(
        "Navigation",
        [
            "📊 Dashboard",
            "📥 Ingest Materials",
            "📁 Manage Materials",
            "🔍 View Chunks",
            "🔍 Chunking Audit",
            "🕸️ Context Graph",
            "🔧 Generate Roadmap",
            "👥 Format by Persona",
            "💬 Ask Your Roadmap",
            "📝 Open Questions",
            "🏗️ Architecture Alignment",
            "🎯 Competitive Intelligence",
            "⚙️ Settings"
        ]
    )

    st.sidebar.markdown("---")
    st.sidebar.caption("Built with Streamlit & Claude")

    # Route to pages
    if page == "📊 Dashboard":
        page_dashboard()
    elif page == "📥 Ingest Materials":
        page_ingest()
    elif page == "📁 Manage Materials":
        page_manage()
    elif page == "🔍 View Chunks":
        page_chunks()
    elif page == "🔍 Chunking Audit":
        page_chunking_audit()
    elif page == "🕸️ Context Graph":
        page_context_graph()
    elif page == "🔧 Generate Roadmap":
        page_generate()
    elif page == "👥 Format by Persona":
        page_format()
    elif page == "💬 Ask Your Roadmap":
        page_ask()
    elif page == "📝 Open Questions":
        page_open_questions()
    elif page == "🏗️ Architecture Alignment":
        page_architecture_alignment()
    elif page == "🎯 Competitive Intelligence":
        page_competitive_intelligence()
    elif page == "⚙️ Settings":
        page_settings()


if __name__ == "__main__":
    main()
